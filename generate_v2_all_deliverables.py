# -*- coding: utf-8 -*-
import os, sys, shutil
sys.stdout.reconfigure(encoding='utf-8')
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE

base_dir = os.path.dirname(os.path.abspath(__file__))
v2_dir = os.path.join(base_dir, "التقرير_الاحترافي_المطور")
output_docx_v2 = os.path.join(v2_dir, "تقرير_جمعية_طبيبي_النصف_سنوي_٢٠٢٦_النسخة_التنفيذية.docx")
output_pptx_v2 = os.path.join(v2_dir, "عرض_تقديمي_جمعية_طبيبي_٢٠٢٦_النسخة_التنفيذية.pptx")
output_slides_v2 = os.path.join(v2_dir, "presentation.html")

print("--- Generating Word V2 ---")
# 1. GENERATE WORD V2
doc = Document()
for section in doc.sections:
    section.top_margin = Inches(0.85)
    section.bottom_margin = Inches(0.85)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)
    section._sectPr.append(parse_xml(f'<w:bidi {nsdecls("w")}/>'))

COLOR_PRIMARY = RGBColor(84, 18, 40)       # Deep Burgundy
COLOR_SECONDARY = RGBColor(201, 169, 110) # Gold
COLOR_DARK = RGBColor(36, 34, 32)
COLOR_MUTED = RGBColor(107, 104, 100)
HEX_PRIMARY = "541228"
HEX_SECONDARY = "C9A96E"
HEX_BG_ALT = "F8F6F0"
HEX_TOTAL = "F2ECE1"

def set_p_rtl(p):
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p._p.get_or_add_pPr().append(parse_xml(f'<w:bidi {nsdecls("w")}/>'))

def set_run_rtl(r, font_name='Cairo', size_pt=11, bold=False, color_rgb=None, italic=False):
    r.font.name = font_name
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.italic = italic
    if color_rgb:
        r.font.color.rgb = color_rgb
    rPr = r._r.get_or_add_rPr()
    rPr.append(parse_xml(f'<w:rtl {nsdecls("w")}/>'))
    rPr.append(parse_xml(f'<w:rFonts {nsdecls("w")} w:ascii="{font_name}" w:hAnsi="{font_name}" w:cs="{font_name}"/>'))

def set_cell_background(cell, hex_color):
    cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:shd {nsdecls("w")} w:val="clear" w:color="auto" w:fill="{hex_color}"/>'))

def set_cell_margins(cell, top=140, bottom=140, left=180, right=180):
    cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:tcMar {nsdecls("w")}><w:top w:w="{top}" w:type="dxa"/><w:bottom w:w="{bottom}" w:type="dxa"/><w:left w:w="{left}" w:type="dxa"/><w:right w:w="{right}" w:type="dxa"/></w:tcMar>'))

def add_heading_1(text):
    p = doc.add_paragraph()
    set_p_rtl(p)
    p.paragraph_format.space_before = Pt(22)
    p.paragraph_format.space_after = Pt(10)
    r = p.add_run(text)
    set_run_rtl(r, font_name='Cairo', size_pt=16, bold=True, color_rgb=COLOR_PRIMARY)
    return p

def add_heading_2(text):
    p = doc.add_paragraph()
    set_p_rtl(p)
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(text)
    set_run_rtl(r, font_name='Cairo', size_pt=13, bold=True, color_rgb=COLOR_PRIMARY)
    return p

def add_body_p(text, bold_prefix=None, italic=False):
    p = doc.add_paragraph()
    set_p_rtl(p)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = 1.25
    if bold_prefix:
        r_b = p.add_run(bold_prefix)
        set_run_rtl(r_b, font_name='Cairo', size_pt=11, bold=True, color_rgb=COLOR_DARK)
    r = p.add_run(text)
    set_rudef add_callout(text, title=None):
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    tbl._tbl.tblPr.append(parse_xml(f'<w:bidiVisual {nsdecls("w")}/>'))
    cell = tbl.cell(0, 0)
    set_cell_background(cell, "F9F7F2")
    set_cell_margins(cell, top=160, bottom=160, left=200, right=200)
    cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="none"/><w:left w:val="none"/><w:bottom w:val="none"/><w:right w:val="single" w:sz="24" w:space="0" w:color="{HEX_PRIMARY}"/></w:tcBorders>'))
    
    p = cell.paragraphs[0]
    set_p_rtl(p)
    if title:
        rt = p.add_run(title + "\n")
        set_run_rtl(rt, font_name='Cairo', size_pt=11, bold=True, color_rgb=COLOR_PRIMARY)
    r = p.add_run(text)
    set_run_rtl(r, font_name='Cairo', size_pt=10.5, italic=True)
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

def style_table(tbl, col_widths, headers, rows_data, is_total_row_present=False):
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    tbl._tbl.tblPr.append(parse_xml(f'<w:bidiVisual {nsdecls("w")}/>'))
    tbl._tbl.tblPr.append(parse_xml(f'<w:tblBorders {nsdecls("w")}><w:top w:val="single" w:sz="4" w:space="0" w:color="E0DDD5"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="E0DDD5"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="EFECE6"/><w:left w:val="none"/><w:right w:val="none"/><w:insideV w:val="none"/></w:tblBorders>'))
    
    hdr_cells = tbl.rows[0].cells
    for i, title in enumerate(headers):
        hdr_cells[i].text = ""
        set_cell_background(hdr_cells[i], HEX_PRIMARY)
        set_cell_margins(hdr_cells[i], top=130, bottom=130, left=150, right=150)
        p = hdr_cells[i].paragraphs[0]
        set_p_rtl(p)
        if title in ["م", "العدد", "رقم الصفحة", "نسبة النمو", "نسبة التنفيذ", "نسبة التغير", "الوزن النسبي", "النسبة", "التاريخ"]:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        else:
            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r = p.add_run(title)
        set_run_rtl(r, font_name='Cairo', size_pt=10, bold=True, color_rgb=RGBColor(255, 255, 255))
            
    for r_idx, row in enumerate(rows_data):
        row_cells = tbl.add_row().cells
        is_total = is_total_row_present and (r_idx == len(rows_data) - 1)
        bg = HEX_TOTAL if is_total else (HEX_BG_ALT if r_idx % 2 == 1 else "FFFFFF")
        
        for c_idx, val in enumerate(row):
            row_cells[c_idx].text = ""
            set_cell_background(row_cells[c_idx], bg)
            set_cell_margins(row_cells[c_idx], top=110, bottom=110, left=130, right=130)
            p = row_cells[c_idx].paragraphs[0]
            set_p_rtl(p)
            
            # Align center for serials, numbers, dates, short percentages
            val_str = str(val).strip()
            if headers[c_idx] in ["م", "العدد", "رقم الصفحة", "نسبة النمو", "نسبة التنفيذ", "نسبة التغير", "الوزن النسبي", "النسبة", "التاريخ"]:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            else:
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                
            r = p.add_run(val_str)
            if is_total:
                set_run_rtl(r, font_name='Cairo', size_pt=9.5, bold=True, color_rgb=COLOR_PRIMARY)
            elif c_idx == 0 and headers[0] == "م":
                set_run_rtl(r, font_name='Cairo', size_pt=9.5, bold=True, color_rgb=COLOR_PRIMARY)
            elif c_idx == 0:
                set_run_rtl(r, font_name='Cairo', size_pt=9.5, bold=True, color_rgb=COLOR_DARK)
            else:
                set_run_rtl(r, font_name='Cairo', size_pt=9.5, bold=False, color_rgb=COLOR_DARK)

    for row in tbl.rows:
        for c_idx, w in enumerate(col_widths):
            row.cells[c_idx].width = Inches(w)
            
    doc.add_paragraph().paragraph_format.space_after = Pt(6)

# Cover
p_cov = doc.add_paragraph()
p_cov.alignment = WD_ALIGN_PARAGRAPH.CENTER
p_cov.paragraph_format.space_before = Pt(70)

r1 = p_cov.add_run("المملكة العربية السعودية | المركز الوطني لتنمية القطاع غير الربحي (ترخيص: ١٠٠٠٧٣٠٧٠٠)\n")
set_run_rtl(r1, font_name='Cairo', size_pt=11, color_rgb=COLOR_MUTED)

r2 = p_cov.add_run("\nجمعية طبيبي الأهلية بالمدينة المنورة\n")
set_run_rtl(r2, font_name='Cairo', size_pt=22, bold=True, color_rgb=COLOR_PRIMARY)

r3 = p_cov.add_run("التقرير النصف سنوي الشامل لعام ٢٠٢٦م\n")
set_run_rtl(r3, font_name='Cairo', size_pt=26, bold=True, color_rgb=COLOR_PRIMARY)

r4 = p_cov.add_run("الفترة من ١ يناير حتى ٣٠ يونيو ٢٠٢٦م\n\n")
set_run_rtl(r4, font_name='Cairo', size_pt=14, bold=True, color_rgb=COLOR_SECONDARY)

r5 = p_cov.add_run("« ثـقـة  •  أثــر  •  اسـتـدامـة »\n\n\n\n")
set_run_rtl(r5, font_name='Cairo', size_pt=18, bold=True, color_rgb=COLOR_PRIMARY)

p_bot = doc.add_paragraph()
p_bot.alignment = WD_ALIGN_PARAGRAPH.CENTER
r6 = p_bot.add_run("إعداد والإشراف التنفيذي:\nأ. بيان بن سعد المحمدي (المدير التنفيذي)\nالمدينة المنورة - حي الفتح | هاتف: (00966555606347) | البريد: (tabibi2025med@gmail.com)")
set_run_rtl(r6, font_name='Cairo', size_pt=10.5, color_rgb=COLOR_DARK)

doc.add_page_break()

# TOC
add_heading_1("فهرس محتويات التقرير")
toc_headers = ["م", "الموضوع / المحور", "المحتوى التفصيلي", "رقم الصفحة"]
toc_data = [
    ["١", "القيادة الرشيدة والرؤية الوطنية", "كلمات خادم الحرمين الشريفين، سمو ولي العهد، وأمير المنطقة", "ص ٣"],
    ["٢", "كلمة رئيس مجلس الإدارة", "رسالة أ.د. منصور محمد النزهة وتوجيهات المجلس (٩ أعضاء)", "ص ٤"],
    ["٣", "الملخص التنفيذي ومصفوفة الـ (KPIs)", "مصفوفة مؤشرات الأداء الـ (١٣) الشاملة ونسب النمو", "ص ٥"],
    ["٤", "الأداء المالي ومصادر الدخل", "مقارنة الإيرادات (H1 2026 vs H1 2025) ونسب النمو", "ص ٦"],
    ["٥", "الموازنة التشغيلية وهيكل السيولة", "نسب تنفيذ الموازنة التقديرية وتوزيع الأرصدة المصرفية", "ص ٧"],
    ["٦", "البيان التفصيلي للمصروفات", "جدول المصروفات التشغيلية بجميع بنوده الـ (١٦)", "ص ٨"],
    ["٧", "البرامج الطبية ورعاية المرضى", "برنامج (جودة حياة) وتفاصيل الحالات الـ (٧) الحرجة المدعومة", "ص ٩"],
    ["٨", "تحليل الحالات المرفوضة", "دراسة أسباب عدم الصرف لـ (١٤) حالة وتوصيات اللائحة", "ص ١٠"],
    ["٩", "شبكة الشراكات الصحية", "اتفاقيات التعاون مع المستشفيات والجهات الصحية الـ (٩)", "ص ١١"],
    ["١٠", "الموارد البشرية والتدريب", "الهيكل الوظيفي، التوطين (١٠٠٪)، والدورات التدريبية الـ (٨)", "ص ١٢"],
    ["١١", "الحوكمة والتحول المؤسسي", "تطبيق نظام قيود، وفر المقر (٢٥ ألف ريال)، ومراحل الخطة", "ص ١٣"],
    ["١٢", "الملاحق الرسمية الكاملة", "الملحق (١): بيان الداعمين الـ (٢٢) كاملاً (٥٨٢,١٦٧.٥٢ ريال)", "ص ١٤"],
    ["١٣", "ملحق الأصول والتجهيزات", "الملحق (٢): بيان الأصول والتجهيزات الثابتة المشتراة لعام ٢٠٢٦م", "ص ١٥"],
    ["١٤", "الخاتمة وقنوات التواصل", "الجهات الإشرافية، المنصات المعتمدة، وبيانات الاتصال", "ص ١٦"]
]
tbl_toc = doc.add_table(rows=1, cols=4)
style_table(tbl_toc, [0.5, 2.2, 3.3, 0.8], toc_headers, toc_data)
doc.add_page_break()

# Royal Leadership
add_heading_1("القيادة الرشيدة والرؤية الوطنية")

tbl_royal = doc.add_table(rows=1, cols=3)
tbl_royal.alignment = WD_TABLE_ALIGNMENT.CENTER
tbl_royal._tbl.tblPr.append(parse_xml(f'<w:bidiVisual {nsdecls("w")}/>'))

royal_cards_data = [
    ("king_salman.jpg", "خادم الحرمين الشريفين\nالملك سلمان بن عبدالعزيز آل سعود", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، وما نراه من مؤسسات خيرية في مختلف المجالات… إلا جانبًا من الجوانب المشرقة لبلادنا.»"),
    ("crown_prince.jpg", "صاحب السمو الملكي\nالأمير محمد بن سلمان بن عبدالعزيز", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة والثقافة والمجالات البحثية، وسنعتمد عليه بشكل رئيسي.»"),
    ("prince_salman.jpg", "صاحب السمو الملكي\nالأمير سلمان بن سلطان بن عبدالعزيز", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية على مستوى المنطقة باعتبارها شريكًا استراتيجيًا للقطاعين العام والخاص في تحسين جودة الحياة وتعزيز الاستقرار الاجتماعي والاقتصادي.»")
]

for col_i, (img_f, title_text, quote_text) in enumerate(royal_cards_data):
    cell = tbl_royal.cell(0, col_i)
    set_cell_background(cell, "F9F7F2")
    set_cell_margins(cell, top=140, bottom=140, left=120, right=120)
    cell._tc.get_or_add_tcPr().append(parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="6" w:space="0" w:color="{HEX_SECONDARY}"/><w:left w:val="single" w:sz="6" w:space="0" w:color="{HEX_SECONDARY}"/><w:bottom w:val="single" w:sz="6" w:space="0" w:color="{HEX_SECONDARY}"/><w:right w:val="single" w:sz="6" w:space="0" w:color="{HEX_SECONDARY}"/></w:tcBorders>'))
    
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    img_path = os.path.join(v2_dir, "assets", "images", img_f)
    if os.path.exists(img_path):
        p.add_run().add_picture(img_path, width=Inches(1.4))
    
    p_t = cell.add_paragraph()
    p_t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_t = p_t.add_run("\n" + title_text + "\n")
    set_run_rtl(r_t, font_name='Cairo', size_pt=10, bold=True, color_rgb=COLOR_PRIMARY)
    
    p_q = cell.add_paragraph()
    set_p_rtl(p_q)
    r_q = p_q.add_run(quote_text)
    set_run_rtl(r_q, font_name='Cairo', size_pt=9.0, italic=True, color_rgb=COLOR_DARK)

for row in tbl_royal.rows:
    for c in row.cells:
        c.width = Inches(2.26)

doc.add_paragraph().paragraph_format.space_after = Pt(8)

# Chairman
add_heading_1("كلمة رئيس مجلس الإدارة")
add_body_p("الحمد لله رب العالمين، والصلاة والسلام على نبينا محمد وعلى آله وصحبه أجمعين.. وبعد:")
add_body_p("يسرني أن أضع بين أيديكم التقرير النصف سنوي لجمعية طبيبي الأهلية، والذي يعكس ما تحقق خلال النصف الأول من عام ٢٠٢٦م من نمو مالي وتشغيلي، وتطور في البنية المؤسسية والحوكمة، وتوسع في الخدمات المقدمة للمستفيدين المرضى في طيبة الطيبة.")
add_body_p("وما تحقق من إنجازات - بعد توفيق الله - هو ثمرة تكامل جهود مجلس الإدارة والجمعية العمومية والإدارة التنفيذية والعاملين والمتطوعين، ودعم الشركاء والمانحين الأفاضل الذين نعتز بثقتهم وإسهامهم في رسالة الجمعية التنموية والإنسانية.")
add_body_p("وننظر إلى هذا التقرير بوصفه أداة للتقييم والتطوير المستمر، لا مجرد عرض للمنجزات؛ بما يساعد على تحديد أولويات المرحلة القادمة، وتعزيز الاستدامة المالية، ورفع الأثر الصحي والاجتماعي المحقق للمستفيدين.")

p_sig = doc.add_paragraph()
set_p_rtl(p_sig)
p_sig.alignment = WD_ALIGN_PARAGRAPH.LEFT
r_sig = p_sig.add_run("أ.د. منصور محمد النزهة\nرئيس مجلس الإدارة | جمعية طبيبي الأهلية")
set_run_rtl(r_sig, font_name='Cairo', size_pt=11, bold=True, color_rgb=COLOR_PRIMARY)

add_heading_2("مجلس الإدارة (٩ أعضاء)")
add_body_p("يقود الجمعية مجلس إدارة مكون من (٩) أعضاء ذوي كفاءات طبية وإدارية واجتماعية، يشرفون على لجان الحوكمة واللجنة التنفيذية ولجنة المساعدات الطبية، لضمان أعلى معايير الشفافية والامتثال المؤسسي.")
doc.add_page_break()

# KPIs
add_heading_1("الملخص التنفيذي ومصفوفة مؤشرات الأداء (KPIs)")
add_body_p("سجل النصف الأول من عام ٢٠٢٦م قفزة نوعية تمثلت في نمو الإيرادات بنسبة (+١٩٢٪) ونمو المساعدات الطبية بنسبة (+٩٤٣٪)، مع توطين كامل للكادر بنسبة (١٠٠٪).")
kpi_table_headers = ["المحور", "مؤشر الأداء (KPI)", "المحقق H1 2026", "المستهدف / المقارنة", "التقييم المؤسسي"]
kpi_table_data = [
    ["المالي", "نمو إجمالي الإيرادات", "٥٨٢,١٦٧ ريال", "+١٩٢٪ (مقابل ١٩٩,٤٧٤ ريال)", "أداء ممتاز"],
    ["المالي", "تنفيذ الموازنة السنوية", "٣٥.٥٧٪", "١,٠٦٠,٦٦٦ من ٢,٩٨١,٧٥٠ ريال", "متوسط (نصف سنة)"],
    ["المالي", "نسبة المصروفات الإدارية", "٥٣.٨٪", "المرتكزة في الرواتب والإيجار", "تتطلب ترشيد"],
    ["المالي", "تغطية الاحتياطي النقدي", "١٢ شهراً", "أرصدة مصرفية (١,٠٠١,٧٥٤ ريال)", "استقرار مالي عالٍ"],
    ["المالي", "مؤشر تركز المانحين", "٤٣٪", "متبرع رئيسي (٢٥٠,٠٠٠ ريال)", "مخاطرة تنوع"],
    ["الطبي", "نمو المساعدات العلاجية", "٢٠٨,٦٠٥ ريال", "+٩٤٣٪ (مقابل ٢٠,٠٠٠ ريال)", "نمو قياسي"],
    ["الطبي", "معدل قبول الحالات", "٣٣.٣٪", "٧ حالات مدعومة من (٢١) متقدمة", "مراجعة اللائحة"],
    ["الطبي", "متوسط كلفة المستفيد", "٢٩,٨٠١ ريال", "عمليات أورام وسرطان متقدمة", "رعاية نوعية"],
    ["الأثر والرضا", "مؤشر تحسن المرضى والرضا", "١٠٠٪", "تحسن كامل للحالات الـ (٧)", "أثر مثبت بالوثائق"],
    ["الموارد البشرية", "نسبة التوطين (السعودة)", "١٠٠٪", "٣ موظفين رسميين + متعاون", "امتثال تام"],
    ["الموارد البشرية", "التدريب والتطوير", "٨ دورات", "استفاد منها موظفان بالكادر", "تطوير مستمر"],
    ["تنمية الموارد", "معدل تحويل طلبات المنح", "٧.٤٪", "قبول منحتين (٤٠,٠٠٠ ريال) من (٢٧)", "استكمال الحوكمة"],
    ["الحوكمة", "تحصيل الذمم المدينة", "٠٪", "١٢,٠٠٠ ريال اشتراكات معلقة", "متابعة التحصيل"]
]
tbl_kpi = doc.add_table(rows=1, cols=5)
style_table(tbl_kpi, [1.0, 1.8, 1.2, 1.6, 1.2], kpi_table_headers, kpi_table_data)

# Financial Tables
add_heading_1("الأداء المالي والموازنة التشغيلية")
add_heading_2("١. مقارنة مصادر الدخل (H1 2026 vs H1 2025)")
rev_headers = ["بند الإيراد", "النصف الأول ٢٠٢٦م", "النصف الأول ٢٠٢٥م", "التغير (ريال)", "نسبة النمو"]
rev_data = [
    ["أموال الزكاة", "٧٠,٠٠٠", "٨٠,٠٠٠", "-١٠,٠٠٠", "-١٣٪"],
    ["علاج مقيد (مساعدات)", "٧٥,٠٠٠", "٢٥,٠٠٠", "+٥٠,٠٠٠", "+٢٠٠٪"],
    ["المتجر الإلكتروني", "١٠,٤٦٩", "١٢٤", "+١٠,٣٤٥", "+٨,٣٤٣٪"],
    ["منصة تبرع", "١,٢٠٣", "١٣,٧٨٦", "-١٢,٥٨٣", "-٩١٪"],
    ["تبرعات ودعم عام", "٤٠٧,٤٩٥", "٦٢,٥٦٤", "+٣٤٤,٩٣١", "+٥٥١٪"],
    ["اشتراكات العضوية", "١٨,٠٠٠", "١٨,٠٠٠", "٠", "٠٪"],
    ["الإجمالي العام", "٥٨٢,١٦٧", "١٩٩,٤٧٤", "+٣٨٢,٦٩٣", "+١٩٢٪"]
]
tbl_rev = doc.add_table(rows=1, cols=5)
style_table(tbl_rev, [1.8, 1.3, 1.3, 1.2, 1.2], rev_headers, rev_data, is_total_row_present=True)

add_heading_2("٢. مستوى تنفيذ الموازنة السنوية لعام ٢٠٢٦م")
bud_headers = ["بند الموازنة", "المستهدف السنوي (ريال)", "المحقق H1 2026 (ريال)", "نسبة التنفيذ"]
bud_data = [
    ["التبرعات والدعم (الإيرادات)", "١,٥٢٧,٠٠٠", "٥٨٢,١٦٧", "٣٨.١٣٪"],
    ["المساعدات العلاجية للمرضى", "٧٥٠,٠٠٠", "٢٠٨,٦٠٥", "٢٧.٨١٪"],
    ["الرواتب والأجور", "٤٧٢,٠٠٠", "١٤٤,٤٠٥", "٣٠.٥٩٪"],
    ["المصروفات التشغيلية", "١٤٢,٣٠٠", "١٠٩,٨٦٩", "٧٧.٢١٪"],
    ["شراء الأصول والتجهيزات", "١٩,٤٥٠", "١٥,٦٢١", "٨٠.٣١٪"],
    ["التطوير المالي والإداري", "٤١,٠٠٠", "٠", "٠٪"],
    ["الحملة الإعلامية وتنمية الموارد", "٣٠,٠٠٠", "٠", "٠٪"],
    ["الإجمالي العام", "٢,٩٨١,٧٥٠", "١,٠٦٠,٦٦٦", "٣٥.٥٧٪"]
]
tbl_bud = doc.add_table(rows=1, cols=4)
style_table(tbl_bud, [2.4, 1.6, 1.6, 1.2], bud_headers, bud_data, is_total_row_present=True)
doc.add_page_break()

add_heading_2("٣. البيان التفصيلي للمصروفات التشغيلية (١٦ بنداً)")
exp_headers = ["م", "بند المصروف", "H1 2026 (ريال)", "H1 2025 (ريال)", "نسبة التغير", "الوزن النسبي"]
exp_data = [
    ["١", "الرواتب الأساسية", "١٤٤,٤٠٥", "٤٥,٢٦٤", "+٢١٩٪", "٥٦.٨٪"],
    ["٢", "الإيجار المكتبي", "٦٣,٣٣٣", "٣٥,٠٠٠", "+٨١٪", "٢٤.٩٪"],
    ["٣", "التأمينات الاجتماعية", "١٤,٧٦٨", "٩,٩٨٠", "+٤٨٪", "٥.٨٪"],
    ["٤", "أجور متعاونين", "١٣,٠٠٠", "٩,٠٦٠", "+٤٣٪", "٥.١٪"],
    ["٥", "المحاسب القانوني (رائد الأحمدي)", "٤,٦٠٠", "٠", "—", "١.٨٪"],
    ["٦", "الكهرباء والخدمات", "٣,٨٦٧", "٠", "—", "١.٥٪"],
    ["٧", "تصميم وتطوير الموقع الإلكتروني", "٣,٠٠٠", "٠", "—", "١.٢٪"],
    ["٨", "نقل وتركيب الأصول للمقر الجديد", "٢,٤٣٠", "٠", "—", "١.٠٪"],
    ["٩", "الهاتف والإنترنت", "١,٣١٦", "١,٣٤٢", "-٢٪", "٠.٥٪"],
    ["١٠", "صيانة متنوعة", "١,٠٦٠", "١,٣٩٣", "-٢٤٪", "٠.٤٪"],
    ["١١", "نظافة ومنظفات", "٩٠٠", "٥٣١", "+٦٩٪", "٠.٣٥٪"],
    ["١٢", "طباعة ومطبوعات", "٥٠٨", "٠", "—", "٠.٢٪"],
    ["١٣", "رسوم مصرفية وعمولات", "٣٨٠", "٠", "—", "٠.١٥٪"],
    ["١٤", "ضيافة واستقبال", "٣٧٥", "٥٩٢", "-٣٧٪", "٠.١٥٪"],
    ["١٥", "أحبار طابعات", "١٨٠", "٠", "—", "٠.٠٧٪"],
    ["١٦", "أدوات مكتبية وقرطاسية", "١٥٢", "٣٦٧", "-٥٩٪", "٠.٠٦٪"],
    ["—", "إجمالي المصروفات التشغيلية", "٢٥٤,٢٧٤", "٦٣,٥٣٦", "+٣٠٠٪", "١٠٠٪"]
]
tbl_exp = doc.add_table(rows=1, cols=6)
style_table(tbl_exp, [0.4, 2.3, 1.1, 1.1, 0.95, 0.95], exp_headers, exp_data, is_total_row_present=True)

add_heading_2("٤. هيكل المركز المالي والسيولة النقدية")
add_body_p("إجمالي الأرصدة المصرفية المتوفرة كما في ٣٠ يونيو ٢٠٢٦م بلغت ", bold_prefix="الأرصدة البنكية: ")
add_body_p("• البنك الأهلي السعودي: (٩٣٠,٧٠٢ ريال) | مصرف الراجحي: (٧١,٠٥٢ ريال) بمجموع: (١,٠٠١,٧٥٤ ريال).")
add_body_p("• الأموال المقيدة لبرامج ومشاريع محددة: (٣٦٧,٠٩٣ ريال) وتمثل (٣٦.٧٪) من السيولة.")
add_body_p("• الأموال غير المقيدة (دعم عام وتشغيل): (٦٣٤,٦٦١ ريال) وتمثل (٦٣.٣٪) من السيولة.")
add_body_p("• صافي الأصول التراكمية: (٩٧٢,٧١٣ ريال) مقابل (٨٦٤,٠٤٥ ريال) في بداية الفترة.")

add_callout(
    "تمت تسوية وسداد كامل الالتزامات المرحلة من عام ٢٠٢٥م بقيمة (١٨,٢١١ ريال) شملت إيجار المقر السابق (٥,٨٣٣ ريال)، أتعاب المحاسب القانوني (٤,٦٠٠ ريال)، إقفال الحسابات (٢,٠٠٠ ريال)، ومخصص نهاية الخدمة (٧٧٨ ريال)، مع بقاء التزام قدره (٥,٠٠٠ ريال) لمؤسسة مؤشرات النجاح عن استكمال ملف الحوكمة.",
    "ملاحظة تدقيقية حول الالتزامات والأصول:"
)
doc.add_page_break()

# Medical
add_heading_1("البرامج والخدمات الطبية ورعاية المرضى")
add_body_p("قدم برنامج «جودة حياة» دعماً طبياً استثنائياً بإجمالي (٢٠٨,٦٠٥.٣١ ريال) لـ (٧) حالات حرجة، محققاً نسبة شفاء وتحسن بلغت (١٠٠٪).")
add_heading_2("١. بيان الحالات السبع المدعومة بالتفصيل")
pat_headers = ["م", "اسم المستفيد", "الجهة العلاجية", "التشخيص الطبي", "المبلغ المعتمد (ريال)"]
pat_data = [
    ["١", "فايز أحمد عبدالعزيز", "المستشفى السعودي الألماني", "سرطان الدم (علاج مناعي وكيماوي)", "١٥٠,٠٠٠"],
    ["٢", "زينب عمر علي", "المستشفى السعودي الألماني", "سرطان نخر العظم", "٣٠,٠٠٠"],
    ["٣", "كندفة محمد عتبة", "مدينة الملك سلمان الطبية", "تنويم ورعاية تحت الملاحظة الفائقة", "٧,٠٠٠"],
    ["٤", "شوق حسن الأنور", "المستشفى السعودي الألماني", "منظار جراحي متقدم", "٧,٠٠٠"],
    ["٥", "سامية سليمان محمد", "مستشفى المواساة بالمدينة", "استئصال كتلة ورمية بالصدر", "٦,٣٥٠"],
    ["٦", "زبيدة شمس الدين خاتم", "المستشفى السعودي الألماني", "ورم بالقولون", "٦,٣٣٠.٣١"],
    ["٧", "محمد أحمد الشرفي", "مستشفى المواساة بالمدينة", "أشعة رنين مغناطيسي تخصصية", "١,٩٢٥"],
    ["—", "إجمالي المساعدات الطبية المعتمدة", "—", "—", "٢٠٨,٦٠٥.٣١"]
]
tbl_pat = doc.add_table(rows=1, cols=5)
style_table(tbl_pat, [0.4, 1.8, 1.8, 1.8, 1.0], pat_headers, pat_data, is_total_row_present=True)

add_heading_2("٢. التحليل التدقيقي لـ (١٤) حالة مرفوضة وأسباب عدم الصرف")
rej_headers = ["سبب عدم الصرف", "العدد", "النسبة", "أبرز الحالات", "التوصية الإدارية"]
rej_data = [
    ["انتهاء صلاحية الإقامة", "٧", "٥٠.٠٪", "بسمة هارون، سيد الأمين، فريدة عظيم، عطور عباس، هاجر الصادق، عبدالله دياب، أحمد خير", "تعديل لائحة المساعدات بالتنسيق مع المانحين للحالات الطارئة"],
    ["تغطية كاملة من جمعية أخرى", "٢", "١٤.٣٪", "هديباء الجهني (مياه بيضاء)، علي قايد (شريان تاجي)", "تفعيل الربط الإلكتروني لمنع ازدواجية الدعم"],
    ["أخطاء بالتقرير الطبي / تشخيص", "٢", "١٤.٣٪", "ريم فواز (ورم ليفي)، جوهرة خان (أخطاء تواريخ)", "إرشاد المستفيد لتصحيح التقارير من المستشفى الشريك"],
    ["وجود تأمين طبي ساري", "١", "٧.١٪", "فؤاد لطف محمد (ميلوما متعددة)", "توجيه المريض للاستفادة من وثيقة التأمين المعتمدة"],
    ["انتهاء تأشيرة الزيارة والسفر", "١", "٧.١٪", "حمزة محمد هندية (سكري نوع أول)", "إغلاق الملف لانتفاء شرط الإقامة النظامية"],
    ["مقبولة ولم تستلم التعميد", "١", "٧.١٪", "مزاهر عبدالله الهادي (ضعف نظر)", "متابعة تسليم التعميد وبدء الخطة العلاجية"]
]
tbl_rej = doc.add_table(rows=1, cols=5)
style_table(tbl_rej, [1.6, 0.6, 0.8, 1.9, 1.9], rej_headers, rej_data)

# Partnerships & Governance
add_heading_1("الشراكات، الموارد البشرية، وخطة النصف الثاني")
add_heading_2("١. الشراكات الصحية الـ (٩) المفعّلة")
add_body_p("• المستشفى السعودي الألماني: جراحات الأورام، مناظير، وسرطانات الدم.")
add_body_p("• مستشفى المواساة بالمدينة المنورة: العمليات الجراحية المتقدمة والأشعة المقطعية والرنين.")
add_body_p("• مدينة الملك سلمان الطبية: الرعاية التخصصية والتنويم والعناية الفائقة.")
add_body_p("• مستشفى د. حامد سليمان الأحمدي: جراحات اليوم الواحد والعيادات الاستشارية.")
add_body_p("• مستشفى واد الطبي & مستشفى المدينة الوطني & مستشفى المدينة الطبي العام.")
add_body_p("• شركة مداواة ورعاية الطبية & جمعية جَنَى لتأهيل الفتيات ذوات الإعاقة.")

add_heading_2("٢. الموارد البشرية والحوكمة المؤسسية")
add_body_p("• الكادر الوظيفي الحالي: (٣) موظفين رسميين (المدير التنفيذي: أ. بيان المحمدي، المسؤول المالي والمشاريع: أ. غدير الحربي، السكرتير التنفيذي: أ. طراد سمان) + محاسب متعاون (أ. محمد الحسن) + مسؤول إعلام قيد الترسيم (أ. فيصل الجهني).")
add_body_p("• نسبة التوطين: (١٠٠٪) في كافة الوظائف، مع إنجاز (٨) دورات تدريبية متخصصة للكادر.")
add_body_p("• التحول الرقمي: تطبيق نظام المحاسبة السحابي «قيود»، وتأسيس الأرشفة الإلكترونية الشاملة.")
add_body_p("• ترشيد النفقات: الانتقال لمقر جديد بإيجار (٤٥,٠٠٠ ريال) بدلاً من (٧٠,٠٠٠ ريال) بتوفير (٢٥,٠٠٠ ريال سنوياً).")

add_heading_2("٣. خطة النصف الثاني ومقترح الفريق الاستشاري (٣ مراحل)")
add_body_p("١. المرحلة الأولى (الشهر الأول): استكمال متطلبات معيار الامتثال والحوكمة ورفع درجة الجمعية المعتمدة، وتفعيل منصة نوى للمنح.")
add_body_p("٢. المرحلة الثانية (الشهر الثاني): تنمية الموارد وإطلاق مبادرة «بطاقة طبيبي» للمزايا والخصومات الصحية للمستفيدين.")
add_body_p("٣. المرحلة الثالثة (الشهر الثالث): إعداد العروض الاستثمارية والاستعداد المبكر للربع الأول من عام ٢٠٢٧م للجهات التي أغلقت موازناتها.")
doc.add_page_break()

# Appendices
add_heading_1("الملاحق الرسمية والتفصيلية")
add_heading_2("الملحق (١): بيان الداعمين التفصيلي لعام ٢٠٢٦م (٥٨٢,١٦٧.٥٢ ريال)")
don_headers = ["م", "الجهة الداعمة / المانح", "التاريخ", "المبلغ (ريال)", "مجال الدعم"]
don_data = [
    ["١", "أسامة جعفر إبراهيم فقيه", "٢٢/٠١/٢٠٢٦", "٥٠,٠٠٠", "زكاة مقيدة"],
    ["٢", "مريم حبيب محمود أحمد", "١٣/٠٢/٢٠٢٦", "٢٠,٠٠٠", "زكاة مقيدة"],
    ["٣", "وقف الشيخ نغيمش الأحمدي (رحمه الله)", "١٦/٠٢/٢٠٢٦", "٥٠,٠٠٠", "٣٥,٠٠٠ علاج + ١٥,٠٠٠ عام"],
    ["٤", "أسامة عدنان حبيب محمود أحمد", "١٦/٠٢/٢٠٢٦", "١٠,٠٠٠", "دعم عام"],
    ["٥", "شركة طابة المطورة للتطوير العمراني", "٢٣/٠٢/٢٠٢٦", "٢٠,٠٠٠", "دعم عام"],
    ["٦", "وقف الشيخ عبدالقادر شيبة الحمد", "٢٧/٠٢/٢٠٢٦", "٥٠,٠٠٠", "دعم عام"],
    ["٧", "شركة حسن محمد حجري", "٢٧/٠٢/٢٠٢٦", "٥,٠٠٠", "دعم عام"],
    ["٨", "سمر فتح الرحمن علي", "٠١/٠٣/٢٠٢٦", "٢,٠٠٠", "دعم عام"],
    ["٩", "سعد بن محمد حسين", "٠٣/٠٣/٢٠٢٦", "٢٥٠,٠٠٠", "دعم عام (أكبر متبرع)"],
    ["١٠", "مربا بنت محمد محروس", "٠٤/٠٣/٢٠٢٦", "٢٠٠", "دعم عام"],
    ["١١", "مؤسسة سعيد محمد مكي", "٠٨/٠٣/٢٠٢٦", "٣,٠٠٠", "دعم عام"],
    ["١٢", "ضيف (فاعل خير)", "٠٩/٠٣/٢٠٢٦", "٥,٠٠٠", "دعم عام"],
    ["١٣", "سلطان محمد الفقيهي", "١٠/٠٣/٢٠٢٦", "٣٠,٠٠٠", "دعم عام"],
    ["١٤", "وقف عبدالرحيم عبدالرزاق", "١٢/٠٣/٢٠٢٦", "١٠,٠٠٠", "دعم عام"],
    ["١٥", "مؤسسة سهيلة شيبة الحمد الخيرية", "٠٩/٠٤/٢٠٢٦", "٢٠,٠٠٠", "علاج ومساعدات طبية"],
    ["١٦", "وقف عبدالعزيز عبدالله أبو زيد", "٢٨/٠٦/٢٠٢٦", "٢٠,٠٠٠", "علاج ومساعدات طبية"],
    ["١٧", "المتجر الإلكتروني للجمعية", "متفرقة", "١٠,٤٦٩", "علاج مقيد"],
    ["١٨", "منصة تبرع الوطنية", "٢٣/٠٢/٢٠٢٦", "١,٢٠٢.٨٨", "علاج مقيد"],
    ["١٩", "متفرقات وتبرعات نقدية", "متفرقة", "٣٩٥.٣٠", "دعم عام"],
    ["٢٠", "حوالات مصرفية صغيرة متفرقة", "متفرقة", "٦,٩٠٠.٣٤", "دعم عام"],
    ["٢١", "رسوم اشتراكات العضوية المحصلة", "متفرقة", "١٨,٠٠٠", "دعم عام"],
    ["—", "الإجمالي العام لبيان الداعمين", "—", "٥٨٢,١٦٧.٥٢", "زكاة: ٧٠ ألف | علاج: ٨٦,٦٧٢ | عام: ٤٢٥,٤٩٦"]
]
tbl_don = doc.add_table(rows=1, cols=5)
style_table(tbl_don, [0.4, 2.3, 1.1, 1.2, 1.8], don_headers, don_data, is_total_row_present=True)

add_heading_2("الملحق (٢): بيان الأصول والتجهيزات الثابتة المشتراة لعام ٢٠٢٦م (١٥,٦٢٠.٨٠ ريال)")
ast_headers = ["م", "الأصل / التجهيز", "العدد", "التاريخ", "القيمة (ريال)", "المورد المعتمد"]
ast_data = [
    ["١", "طابعة ليزر ملون HP", "١", "١١/٠٥/٢٠٢٦", "١,٣٥٠", "شركة سمرة الرقمية"],
    ["٢", "مكتب سكرتارية خشب بني", "٦", "٢٧/٠٦/٢٠٢٦", "٤,٦٨٠", "الصفوة الجديدة للأثاث"],
    ["٣", "كرسي دوار جلد رصاصي", "٦", "٢٧/٠٦/٢٠٢٦", "٢,٧٠٠", "الصفوة الجديدة للأثاث"],
    ["٤", "مكيفات أوجين ٢٤ وحدة", "٣", "٢٧/٠٦/٢٠٢٦", "٤,٥٩٠.٨٠", "محل بن بلال للأجهزة"],
    ["٥", "خزينة حديدية للمستندات", "١", "٢٧/٠٦/٢٠٢٦", "١,٢٥٠", "الصفوة الجديدة للأثاث"],
    ["٦", "كرسي انتظار كروم للمراجعين", "٣", "٢٧/٠٦/٢٠٢٦", "١,٠٥٠", "مؤسسة الشرق هوم"],
    ["—", "إجمالي مشتريات الأصول", "—", "—", "١٥,٦٢٠.٨٠", "مقابل ٣٤,٧٧٥.٥٠ ريال في H1 2025"]
]
tbl_ast = doc.add_table(rows=1, cols=6)
style_table(tbl_ast, [0.4, 2.0, 0.6, 1.0, 1.2, 1.6], ast_headers, ast_data, is_total_row_present=True)

doc.save(output_docx_v2)
print(f"Generated Word V2 successfully: {output_docx_v2}")�ى المواساة بالمدينة", "أشعة رنين مغناطيسي تخصصية", "١,٩٢٥"],
    ["—", "إجمالي المساعدات الطبية المعتمدة", "—", "—", "٢٠٨,٦٠٥.٣١"]
]
tbl_pat = doc.add_table(rows=1, cols=5)
style_table(tbl_pat, [0.5, 1.8, 1.8, 1.8, 1.3], pat_headers, pat_data, is_total_row_present=True)

add_heading_2("٢. التحليل التدقيقي لـ (١٤) حالة مرفوضة وأسباب عدم الصرف")
rej_headers = ["سبب عدم الصرف", "العدد", "النسبة", "أبرز الحالات", "التوصية الإدارية"]
rej_data = [
    ["انتهاء صلاحية الإقامة", "٧", "٥٠.٠٪", "بسمة هارون، سيد الأمين، فريدة عظيم، عطور عباس، هاجر الصادق، عبدالله دياب، أحمد خير", "تعديل لائحة المساعدات بالتنسيق مع المانحين للحالات الطارئة"],
    ["تغطية كاملة من جمعية أخرى", "٢", "١٤.٣٪", "هديباء الجهني (مياه بيضاء)، علي قايد (شريان تاجي)", "تفعيل الربط الإلكتروني لمنع ازدواجية الدعم"],
    ["أخطاء بالتقرير الطبي / تشخيص", "٢", "١٤.٣٪", "ريم فواز (ورم ليفي)، جوهرة خان (أخطاء تواريخ)", "إرشاد المستفيد لتصحيح التقارير من المستشفى الشريك"],
    ["وجود تأمين طبي ساري", "١", "٧.١٪", "فؤاد لطف محمد (ميلوما متعددة)", "توجيه المريض للاستفادة من وثيقة التأمين المعتمدة"],
    ["انتهاء تأشيرة الزيارة والسفر", "١", "٧.١٪", "حمزة محمد هندية (سكري نوع أول)", "إغلاق الملف لانتفاء شرط الإقامة النظامية"],
    ["مقبولة ولم تستلم التعميد", "١", "٧.١٪", "مزاهر عبدالله الهادي (ضعف نظر)", "متابعة تسليم التعميد وبدء الخطة العلاجية"]
]
tbl_rej = doc.add_table(rows=1, cols=5)
style_table(tbl_rej, [1.6, 0.6, 0.8, 2.0, 2.0], rej_headers, rej_data)

# Partnerships & Governance
add_heading_1("الشراكات، الموارد البشرية، وخطة النصف الثاني")
add_heading_2("١. الشراكات الصحية الـ (٩) المفعّلة")
add_body_p("• المستشفى السعودي الألماني: جراحات الأورام، مناظير، وسرطانات الدم.")
add_body_p("• مستشفى المواساة بالمدينة المنورة: العمليات الجراحية المتقدمة والأشعة المقطعية والرنين.")
add_body_p("• مدينة الملك سلمان الطبية: الرعاية التخصصية والتنويم والعناية الفائقة.")
add_body_p("• مستشفى د. حامد سليمان الأحمدي: جراحات اليوم الواحد والعيادات الاستشارية.")
add_body_p("• مستشفى واد الطبي & مستشفى المدينة الوطني & مستشفى المدينة الطبي العام.")
add_body_p("• شركة مداواة ورعاية الطبية & جمعية جَنَى لتأهيل الفتيات ذوات الإعاقة.")

add_heading_2("٢. الموارد البشرية والحوكمة المؤسسية")
add_body_p("• الكادر الوظيفي الحالي: (٣) موظفين رسميين (المدير التنفيذي: أ. بيان المحمدي، المسؤول المالي والمشاريع: أ. غدير الحربي، السكرتير التنفيذي: أ. طراد سمان) + محاسب متعاون (أ. محمد الحسن) + مسؤول إعلام قيد الترسيم (أ. فيصل الجهني).")
add_body_p("• نسبة التوطين: (١٠٠٪) في كافة الوظائف، مع إنجاز (٨) دورات تدريبية متخصصة للكادر.")
add_body_p("• التحول الرقمي: تطبيق نظام المحاسبة السحابي «قيود»، وتأسيس الأرشفة الإلكترونية الشاملة.")
add_body_p("• ترشيد النفقات: الانتقال لمقر جديد بإيجار (٤٥,٠٠٠ ريال) بدلاً من (٧٠,٠٠٠ ريال) بتوفير (٢٥,٠٠٠ ريال سنوياً).")

add_heading_2("٣. خطة النصف الثاني ومقترح الفريق الاستشاري (٣ مراحل)")
add_body_p("١. المرحلة الأولى (الشهر الأول): استكمال متطلبات معيار الامتثال والحوكمة ورفع درجة الجمعية المعتمدة، وتفعيل منصة نوى للمنح.")
add_body_p("٢. المرحلة الثانية (الشهر الثاني): تنمية الموارد وإطلاق مبادرة «بطاقة طبيبي» للمزايا والخصومات الصحية للمستفيدين.")
add_body_p("٣. المرحلة الثالثة (الشهر الثالث): إعداد العروض الاستثمارية والاستعداد المبكر للربع الأول من عام ٢٠٢٧م للجهات التي أغلقت موازناتها.")
doc.add_page_break()

# Appendices
add_heading_1("الملاحق الرسمية والتفصيلية")
add_heading_2("الملحق (١): بيان الداعمين التفصيلي لعام ٢٠٢٦م (٥٨٢,١٦٧.٥٢ ريال)")
don_headers = ["م", "الجهة الداعمة / المانح", "التاريخ", "المبلغ (ريال)", "مجال الدعم"]
don_data = [
    ["١", "أسامة جعفر إبراهيم فقيه", "٢٢/٠١/٢٠٢٦", "٥٠,٠٠٠", "زكاة مقيدة"],
    ["٢", "مريم حبيب محمود أحمد", "١٣/٠٢/٢٠٢٦", "٢٠,٠٠٠", "زكاة مقيدة"],
    ["٣", "وقف الشيخ نغيمش الأحمدي (رحمه الله)", "١٦/٠٢/٢٠٢٦", "٥٠,٠٠٠", "٣٥,٠٠٠ علاج + ١٥,٠٠٠ عام"],
    ["٤", "أسامة عدنان حبيب محمود أحمد", "١٦/٠٢/٢٠٢٦", "١٠,٠٠٠", "دعم عام"],
    ["٥", "شركة طابة المطورة للتطوير العمراني", "٢٣/٠٢/٢٠٢٦", "٢٠,٠٠٠", "دعم عام"],
    ["٦", "وقف الشيخ عبدالقادر شيبة الحمد", "٢٧/٠٢/٢٠٢٦", "٥٠,٠٠٠", "دعم عام"],
    ["٧", "شركة حسن محمد حجري", "٢٧/٠٢/٢٠٢٦", "٥,٠٠٠", "دعم عام"],
    ["٨", "سمر فتح الرحمن علي", "٠١/٠٣/٢٠٢٦", "٢,٠٠٠", "دعم عام"],
    ["٩", "سعد بن محمد حسين", "٠٣/٠٣/٢٠٢٦", "٢٥٠,٠٠٠", "دعم عام (أكبر متبرع)"],
    ["١٠", "مربا بنت محمد محروس", "٠٤/٠٣/٢٠٢٦", "٢٠٠", "دعم عام"],
    ["١١", "مؤسسة سعيد محمد مكي", "٠٨/٠٣/٢٠٢٦", "٣,٠٠٠", "دعم عام"],
    ["١٢", "ضيف (فاعل خير)", "٠٩/٠٣/٢٠٢٦", "٥,٠٠٠", "دعم عام"],
    ["١٣", "سلطان محمد الفقيهي", "١٠/٠٣/٢٠٢٦", "٣٠,٠٠٠", "دعم عام"],
    ["١٤", "وقف عبدالرحيم عبدالرزاق", "١٢/٠٣/٢٠٢٦", "١٠,٠٠٠", "دعم عام"],
    ["١٥", "مؤسسة سهيلة شيبة الحمد الخيرية", "٠٩/٠٤/٢٠٢٦", "٢٠,٠٠٠", "علاج ومساعدات طبية"],
    ["١٦", "وقف عبدالعزيز عبدالله أبو زيد", "٢٨/٠٦/٢٠٢٦", "٢٠,٠٠٠", "علاج ومساعدات طبية"],
    ["١٧", "المتجر الإلكتروني للجمعية", "متفرقة", "١٠,٤٦٩", "علاج مقيد"],
    ["١٨", "منصة تبرع الوطنية", "٢٣/٠٢/٢٠٢٦", "١,٢٠٢.٨٨", "علاج مقيد"],
    ["١٩", "متفرقات وتبرعات نقدية", "متفرقة", "٣٩٥.٣٠", "دعم عام"],
    ["٢٠", "حوالات مصرفية صغيرة متفرقة", "متفرقة", "٦,٩٠٠.٣٤", "دعم عام"],
    ["٢١", "رسوم اشتراكات العضوية المحصلة", "متفرقة", "١٨,٠٠٠", "دعم عام"],
    ["—", "الإجمالي العام لبيان الداعمين", "—", "٥٨٢,١٦٧.٥٢", "زكاة: ٧٠ ألف | علاج: ٨٥,٦٧٢ | عام: ٤٢٦,٤٩٦"]
]
tbl_don = doc.add_table(rows=1, cols=5)
style_table(tbl_don, [0.5, 2.3, 1.1, 1.2, 2.1], don_headers, don_data, is_total_row_present=True)

add_heading_2("الملحق (٢): بيان الأصول والتجهيزات الثابتة المشتراة لعام ٢٠٢٦م (١٥,٦٢٠.٨٠ ريال)")
ast_headers = ["الأصل / التجهيز", "العدد", "التاريخ", "القيمة (ريال)", "المورد المعتمد"]
ast_data = [
    ["طابعة ليزر ملون HP", "١", "١١/٠٥/٢٠٢٦", "١,٣٥٠", "شركة سمرة الرقمية"],
    ["مكتب سكرتارية خشب بني", "٦", "٢٧/٠٦/٢٠٢٦", "٤,٦٨٠", "الصفوة الجديدة للأثاث"],
    ["كرسي دوار جلد رصاصي", "٦", "٢٧/٠٦/٢٠٢٦", "٢,٧٠٠", "الصفوة الجديدة للأثاث"],
    ["مكيفات أوجين ٢٤ وحدة", "٣", "٢٧/٠٦/٢٠٢٦", "٤,٥٩٠.٨٠", "محل بن بلال للأجهزة"],
    ["خزينة حديدية للمستندات", "١", "٢٧/٠٦/٢٠٢٦", "١,٢٥٠", "الصفوة الجديدة للأثاث"],
    ["كرسي انتظار كروم للمراجعين", "٣", "٢٧/٠٦/٢٠٢٦", "١,٠٥٠", "مؤسسة الشرق هوم"],
    ["إجمالي مشتريات الأصول", "—", "—", "١٥,٦٢٠.٨٠", "مقابل ٣٤,٧٧٥.٥٠ ريال في H1 2025"]
]
tbl_ast = doc.add_table(rows=1, cols=5)
style_table(tbl_ast, [2.2, 0.6, 1.1, 1.2, 2.1], ast_headers, ast_data, is_total_row_present=True)

doc.save(output_docx_v2)
print(f"Generated Word V2 successfully: {output_docx_v2}")

print("--- Generating PPTX V2 ---")
# 2. GENERATE PPTX V2
prs = Presentation()
prs.slide_width = PInches(13.333)
prs.slide_height = PInches(7.5)

C_PRIMARY = PRGBColor(84, 18, 40)
C_PRIMARY_DARK = PRGBColor(46, 11, 23)
C_SECONDARY = PRGBColor(201, 169, 110)
C_SECONDARY_LT = PRGBColor(223, 202, 155)
C_WHITE = PRGBColor(255, 255, 255)
C_DARK = PRGBColor(36, 34, 32)
C_MUTED = PRGBColor(107, 104, 100)
C_SUCCESS = PRGBColor(30, 130, 76)
C_WARNING = PRGBColor(217, 130, 43)
C_DANGER = PRGBColor(192, 57, 43)

def add_blank_slide_p(bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
    return slide

def add_header_p(slide, title, eyebrow=None, dark_theme=False):
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(0.4), PInches(11.733), PInches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    if eyebrow:
        p_eye = tf.paragraphs[0]
        p_eye.alignment = PP_ALIGN.RIGHT
        r_eye = p_eye.add_run()
        r_eye.text = eyebrow.upper()
        r_eye.font.name = 'Cairo'
        r_eye.font.size = PPt(11)
        r_eye.font.bold = True
        r_eye.font.color.rgb = C_SECONDARY if dark_theme else PRGBColor(166, 133, 71)
        p_title = tf.add_paragraph()
    else:
        p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.RIGHT
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = 'Cairo'
    r_title.font.size = PPt(22)
    r_title.font.bold = True
    r_title.font.color.rgb = C_WHITE if dark_theme else C_PRIMARY

def add_card_p(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(top), PInches(width), PInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = PPt(1.5)
    return shape

# Slide 1: Cover
s1 = add_blank_slide_p(C_PRIMARY_DARK)
tb1 = s1.shapes.add_textbox(PInches(1.0), PInches(1.5), PInches(11.333), PInches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r = p1.add_run()
r.text = "جمعية طبيبي الأهلية بالمدينة المنورة\n"
r.font.name = 'Cairo'
r.font.size = PPt(18)
r.font.color.rgb = C_SECONDARY_LT

p2 = tf1.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r = p2.add_run()
r.text = "التقرير النصف سنوي الشامل لعام ٢٠٢٦م\n"
r.font.name = 'Cairo'
r.font.size = PPt(34)
r.font.bold = True
r.font.color.rgb = C_WHITE

p3 = tf1.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r = p3.add_run()
r.text = "الفترة من ١ يناير حتى ٣٠ يونيو ٢٠٢٦م  |  « ثـقـة  •  أثــر  •  اسـتـدامـة »\n\n"
r.font.name = 'Cairo'
r.font.size = PPt(16)
r.font.color.rgb = C_SECONDARY

tb_ft = s1.shapes.add_textbox(PInches(1.0), PInches(6.0), PInches(11.333), PInches(1.0))
p_ft = tb_ft.text_frame.paragraphs[0]
p_ft.alignment = PP_ALIGN.CENTER
r = p_ft.add_run()
r.text = "ترخيص المركز الوطني لتنمية القطاع غير الربحي رقم: (١٠٠٠٧٣٠٧٠٠) | إشراف: أ. بيان بن سعد المحمدي - المدير التنفيذي"
r.font.name = 'Cairo'
r.font.size = PPt(11)
r.font.color.rgb = PRGBColor(200, 200, 200)

# Slide 2: Royal Leadership
s2 = add_blank_slide_p(PRGBColor(250, 248, 245))
add_header_p(s2, "القيادة الرشيدة والرؤية الوطنية", "الرؤية والتمكين")
quotes = [
    ("crown_prince.jpg", "صاحب السمو الملكي\nالأمير محمد بن سلمان بن عبدالعزيز", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة والثقافة والمجالات البحثية، وسنعتمد عليه بشكل رئيسي.»"),
    ("king_salman.jpg", "خادم الحرمين الشريفين\nالملك سلمان بن عبدالعزيز آل سعود", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، وما نراه من مؤسسات خيرية في مختلف المجالات… إلا جانبًا من الجوانب المشرقة لبلادنا.»"),
    ("prince_salman.jpg", "صاحب السمو الملكي\nالأمير سلمان بن سلطان بن عبدالعزيز", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية على مستوى المنطقة باعتبارها شريكًا استراتيجيًا للقطاعين العام والخاص في تحسين جودة الحياة وتعزيز الاستقرار.»")
]
for idx, (img_f, ldr, qt) in enumerate(quotes):
    left = 0.8 + idx * 3.95
    add_card_p(s2, left, 1.6, 3.8, 5.3, C_PRIMARY_DARK, C_SECONDARY)
    img_path = os.path.join(v2_dir, "assets", "images", img_f)
    if os.path.exists(img_path):
        s2.shapes.add_picture(img_path, PInches(left + 1.15), PInches(1.8), PInches(1.5), PInches(1.8))
    tb = s2.shapes.add_textbox(PInches(left + 0.2), PInches(3.7), PInches(3.4), PInches(3.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ldr + "\n"
    r.font.name = 'Cairo'
    r.font.size = PPt(11)
    r.font.bold = True
    r.font.color.rgb = C_SECONDARY
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r = p2.add_run()
    r.text = qt
    r.font.name = 'Cairo'
    r.font.size = PPt(9.5)
    r.font.italic = True
    r.font.color.rgb = C_WHITE

# Slide 3: KPIs Dashboard
s3 = add_blank_slide_p(PRGBColor(250, 248, 245))
add_header_p(s3, "مصفوفة مؤشرات الأداء الرئيسية (KPIs Dashboard)", "لوحة التحكم")
kpis = [
    ("نمو الإيرادات", "+١٩٢٪", "٥٨٢,١٦٧ ريال مقابل ١٩٩ ألف في ٢٠٢٥", C_SUCCESS),
    ("المساعدات الطبية", "+٩٤٣٪", "٢٠٨,٦٠٥ ريال لـ ٧ حالات حرجة", C_SUCCESS),
    ("تنفيذ الموازنة", "٣٥.٥٧٪", "١,٠٦٠,٦٦٦ من ٢,٩٨١,٧٥٠ ريال معتمد", C_WARNING),
    ("الاحتياطي النقدي", "١٢ شهراً", "أرصدة ١,٠٠١,٧٥٤ ريال تغطي التشغيل", C_SUCCESS),
    ("نسبة التوطين", "١٠٠٪", "٣ موظفين رسميين + محاسب متعاون", C_SUCCESS),
    ("قبول الحالات", "٣٣.٣٪", "٧ حالات مقبولة من أصل ٢١ متقدمة", C_WARNING),
    ("متوسط كلفة المريض", "٢٩,٨٠١ ر.س", "تغطية جراحات وأورام معقدة", C_PRIMARY),
    ("نسبة تحسن الحالات", "١٠٠٪", "تحسن كامل لكافة المرضى المدعومين", C_SUCCESS)
]
for idx, (kt, kv, kd, kc) in enumerate(kpis):
    r_i = idx // 4
    c_i = idx % 4
    left = 0.8 + c_i * 2.95
    top = 1.8 + r_i * 2.5
    add_card_p(s3, left, top, 2.8, 2.3, C_WHITE, PRGBColor(225, 220, 210))
    tb = s3.shapes.add_textbox(PInches(left + 0.15), PInches(top + 0.15), PInches(2.5), PInches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = kt
    r.font.name = 'Cairo'
    r.font.size = PPt(11)
    r.font.bold = True
    r.font.color.rgb = C_MUTED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r = p2.add_run()
    r.text = kv
    r.font.name = 'Cairo'
    r.font.size = PPt(22)
    r.font.bold = True
    r.font.color.rgb = kc
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.RIGHT
    r = p3.add_run()
    r.text = kd
    r.font.name = 'Cairo'
    r.font.size = PPt(9.5)
    r.font.color.rgb = C_DARK

# Slide 4: Closing
s4 = add_blank_slide_p(C_PRIMARY_DARK)
tb_e = s4.shapes.add_textbox(PInches(1.0), PInches(2.0), PInches(11.333), PInches(4.0))
tf_e = tb_e.text_frame
tf_e.word_wrap = True
p = tf_e.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "شكراً لثقتكم ودعمكم المستمر لرسالة جمعية طبيبي الأهلية\n\n"
r.font.name = 'Cairo'
r.font.size = PPt(26)
r.font.bold = True
r.font.color.rgb = C_SECONDARY
p2 = tf_e.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r = p2.add_run()
r.text = "هاتف: (00966555606347)  |  البريد: (tabibi2025med@gmail.com)  |  المدينة المنورة - حي الفتح\nالمركز الوطني لتنمية القطاع غير الربحي - ترخيص: (١٠٠٠٧٣٠٧٠٠)"
r.font.name = 'Cairo'
r.font.size = PPt(12)
r.font.color.rgb = C_SECONDARY_LT

prs.save(output_pptx_v2)
print(f"Generated PPTX V2 successfully: {output_pptx_v2}")

print("--- Generating Web Slide Deck V2 ---")
# 3. GENERATE SLIDE DECK V2
src_pres = os.path.join(base_dir, "التقرير_الجديد", "presentation.html")
if os.path.exists(src_pres):
    shutil.copy(src_pres, output_slides_v2)
    print(f"Generated Web Slide Deck V2 successfully: {output_slides_v2}")

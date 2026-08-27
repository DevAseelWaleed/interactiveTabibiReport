# -*- coding: utf-8 -*-
import os, sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

base_dir = os.path.dirname(os.path.abspath(__file__))
output_pptx = os.path.join(base_dir, "التقرير_الجديد", "عرض_تقديمي_جمعية_طبيبي_النصف_سنوي_٢٠٢٦.pptx")

prs = Presentation()
# Set 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Constants
C_PRIMARY = RGBColor(107, 29, 58)      # #6B1D3A Dark Burgundy
C_PRIMARY_DARK = RGBColor(46, 11, 23) # #2E0B17
C_SECONDARY = RGBColor(201, 169, 110) # #C9A96E Gold
C_SECONDARY_LT = RGBColor(223, 202, 155)
C_WHITE = RGBColor(255, 255, 255)
C_DARK = RGBColor(36, 34, 32)
C_MUTED = RGBColor(107, 104, 100)
C_BG_CARD = RGBColor(255, 255, 255)
C_BG_ALT = RGBColor(241, 239, 234)
C_SUCCESS = RGBColor(30, 130, 76)
C_WARNING = RGBColor(217, 130, 43)
C_DANGER = RGBColor(192, 57, 43)

def add_blank_slide(bg_color=None):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide

def add_header(slide, title, eyebrow=None, dark_theme=False):
    # Header bar
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    if eyebrow:
        p_eye = tf.paragraphs[0]
        p_eye.alignment = PP_ALIGN.RIGHT
        r_eye = p_eye.add_run()
        r_eye.text = eyebrow.upper()
        r_eye.font.name = 'Cairo'
        r_eye.font.size = Pt(11)
        r_eye.font.bold = True
        r_eye.font.color.rgb = C_SECONDARY if dark_theme else RGBColor(166, 133, 71)
        
        p_title = tf.add_paragraph()
    else:
        p_title = tf.paragraphs[0]
        
    p_title.alignment = PP_ALIGN.RIGHT
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = 'Cairo'
    r_title.font.size = Pt(22)
    r_title.font.bold = True
    r_title.font.color.rgb = C_WHITE if dark_theme else C_PRIMARY

def add_card(slide, left, top, width, height, bg_color=C_BG_CARD, border_color=C_SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# =============================================================================
# SLIDE 1: COVER SLIDE
# =============================================================================
s1 = add_blank_slide(C_PRIMARY_DARK)

# Decorative gold banner
top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
top_band.fill.solid()
top_band.fill.fore_color.rgb = C_SECONDARY
top_band.line.fill.background()

# Center text
tb_cover = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
tf1 = tb_cover.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "جمعية طبيبي الأهلية بالمدينة المنورة\n"
r1.font.name = 'Cairo'
r1.font.size = Pt(18)
r1.font.color.rgb = C_SECONDARY_LT

p2 = tf1.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "التقرير النصف سنوي الشامل لعام ٢٠٢٦م\n"
r2.font.name = 'Cairo'
r2.font.size = Pt(36)
r2.font.bold = True
r2.font.color.rgb = C_WHITE

p3 = tf1.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "الفترة من ١ يناير إلى ٣٠ يونيو ٢٠٢٦م\n\n"
r3.font.name = 'Cairo'
r3.font.size = Pt(16)
r3.font.color.rgb = C_SECONDARY

p4 = tf1.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run()
r4.text = "« ثـقـة  •  أثــر  •  اسـتـدامـة »\n"
r4.font.name = 'Cairo'
r4.font.size = Pt(22)
r4.font.bold = True
r4.font.color.rgb = C_SECONDARY_LT

# Footer info
tb_cov_ft = s1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(1.0))
tf_cft = tb_cov_ft.text_frame
p_cft = tf_cft.paragraphs[0]
p_cft.alignment = PP_ALIGN.CENTER
r_cft = p_cft.add_run()
r_cft.text = "ترخيص المركز الوطني لتنمية القطاع غير الربحي رقم: ١٠٠٠٧٣٠٧٠٠ | إشراف: أ. بيان بن سعد المحمدي - المدير التنفيذي"
r_cft.font.name = 'Cairo'
r_cft.font.size = Pt(11)
r_cft.font.color.rgb = RGBColor(200, 200, 200)

# =============================================================================
# SLIDE 2: ROYAL LEADERSHIP
# =============================================================================
s2 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s2, "القيادة الرشيدة والرؤية الوطنية", "الرؤية والتمكين")

quotes = [
    ("خادم الحرمين الشريفين\nالملك سلمان بن عبدالعزيز آل سعود", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، وما نراه من مؤسسات خيرية في مختلف المجالات… إلا جانبًا من الجوانب المشرقة لبلادنا.»"),
    ("صاحب السمو الملكي\nالأمير محمد بن سلمان بن عبدالعزيز", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة والثقافة والمجالات البحثية، وسنعتمد عليه بشكل رئيسي.»"),
    ("صاحب السمو الملكي\nالأمير سلمان بن سلطان بن عبدالعزيز", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية على مستوى المنطقة باعتبارها شريكًا استراتيجيًا للقطاعين العام والخاص في تحسين جودة الحياة وتعزيز الاستقرار.»")
]

for idx, (leader, quote) in enumerate(quotes):
    left = 0.8 + idx * 3.95
    card = add_card(s2, left, 1.8, 3.8, 5.0, C_PRIMARY_DARK, C_SECONDARY)
    tb = s2.shapes.add_textbox(Inches(left + 0.2), Inches(2.0), Inches(3.4), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_l = tf.paragraphs[0]
    p_l.alignment = PP_ALIGN.CENTER
    r_l = p_l.add_run()
    r_l.text = leader + "\n\n"
    r_l.font.name = 'Cairo'
    r_l.font.size = Pt(13)
    r_l.font.bold = True
    r_l.font.color.rgb = C_SECONDARY
    
    p_q = tf.add_paragraph()
    p_q.alignment = PP_ALIGN.RIGHT
    r_q = p_q.add_run()
    r_q.text = quote
    r_q.font.name = 'Cairo'
    r_q.font.size = Pt(11)
    r_q.font.italic = True
    r_q.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 3: CHAIRMAN ADDRESS & HIGHLIGHTS
# =============================================================================
s3 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s3, "كلمة رئيس مجلس الإدارة والملخص التنفيذي", "القيادة المؤسسية")

# Speech card (Left/Main)
add_card(s3, 0.8, 1.8, 7.5, 5.0, C_WHITE, RGBColor(230, 225, 215))
tb_sp = s3.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(6.9), Inches(4.5))
tf_sp = tb_sp.text_frame
tf_sp.word_wrap = True
p_sp = tf_sp.paragraphs[0]
p_sp.alignment = PP_ALIGN.RIGHT
r_sp = p_sp.add_run()
r_sp.text = "يسرني أن أضع بين أيديكم التقرير النصف سنوي لجمعية طبيبي الأهلية، والذي يعكس ما تحقق خلال النصف الأول من عام ٢٠٢٦م من نمو مالي وتشغيلي، وتطور في البنية المؤسسية والحوكمة، وتوسع في الخدمات المقدمة للمستفيدين المرضى في طيبة الطيبة.\n\nوما تحقق من إنجازات - بعد توفيق الله - هو ثمرة تكامل جهود مجلس الإدارة والجمعية العمومية والإدارة التنفيذية والعاملين والمتطوعين، ودعم الشركاء والمانحين الأفاضل الذين نعتز بثقتهم وإسهامهم.\n\nأ.د. منصور محمد النزهة - رئيس مجلس الإدارة"
r_sp.font.name = 'Cairo'
r_sp.font.size = Pt(12)
r_sp.font.color.rgb = C_DARK

# Highlights Summary on the Right
add_card(s3, 8.6, 1.8, 3.9, 5.0, C_PRIMARY, C_SECONDARY)
tb_hi = s3.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(3.5), Inches(4.5))
tf_hi = tb_hi.text_frame
tf_hi.word_wrap = True

p_hit = tf_hi.paragraphs[0]
p_hit.alignment = PP_ALIGN.RIGHT
r_hit = p_hit.add_run()
r_hit.text = "أبرز أرقام النصف الأول:\n\n"
r_hit.font.name = 'Cairo'
r_hit.font.size = Pt(14)
r_hit.font.bold = True
r_hit.font.color.rgb = C_SECONDARY

items_hi = [
    "• إيرادات محققة: ٥٨٢,١٦٧ ريال (+١٩٢٪)",
    "• مساعدات علاجية: ٢٠٨,٦٠٥ ريال (+٩٤٣٪)",
    "• ٧ حالات حرجة تم دعمها بنجاح ١٠٠٪",
    "• ٩ شراكات صحية استراتيجية مفعّلة",
    "• نسبة التوطين بالكادر: ١٠٠٪",
    "• سيولة وأرصدة بنكية: ١,٠٠١,٧٥٤ ريال"
]
for it in items_hi:
    p_it = tf_hi.add_paragraph()
    p_it.alignment = PP_ALIGN.RIGHT
    r_it = p_it.add_run()
    r_it.text = it + "\n"
    r_it.font.name = 'Cairo'
    r_it.font.size = Pt(11)
    r_it.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 4: MASTER KPIS DASHBOARD
# =============================================================================
s4 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s4, "لوحة مؤشرات الأداء الرئيسية (Master KPIs)", "مصفوفة القياس التراكمي")

kpi_cards = [
    ("نمو الإيرادات", "+١٩٢٪", "٥٨٢,١٦٧ ريال مقارنة بـ ١٩٩ ألف في ٢٠٢٥", C_SUCCESS),
    ("المساعدات الطبية", "+٩٤٣٪", "٢٠٨,٦٠٥ ريال لـ ٧ حالات حرجة", C_SUCCESS),
    ("تنفيذ الموازنة", "٣٥.٥٧٪", "١,٠٦٠,٦٦٦ من ٢,٩٨١,٧٥٠ ريال معتمد", C_WARNING),
    ("الاحتياطي النقدي", "١٢ شهراً", "أرصدة ١,٠٠١,٧٥٤ ريال تغطي النفقات التشغيلية", C_SUCCESS),
    ("نسبة التوطين", "١٠٠٪", "٣ موظفين رسميين + محاسب متعاون", C_SUCCESS),
    ("معدل قبول الحالات", "٣٣.٣٪", "٧ حالات مقبولة من أصل ٢١ متقدمة", C_WARNING),
    ("متوسط كلفة المريض", "٢٩,٨٠١ ر.س", "تغطية عمليات أورام وسرطان متقدمة", C_PRIMARY),
    ("نسبة تحسن الحالات", "١٠٠٪", "تحسن كامل لكافة المرضى المدعومين", C_SUCCESS)
]

for idx, (title, val, desc, col) in enumerate(kpi_cards):
    r_idx = idx // 4
    c_idx = idx % 4
    left = 0.8 + c_idx * 2.95
    top = 1.8 + r_idx * 2.5
    
    add_card(s4, left, top, 2.8, 2.3, C_WHITE, RGBColor(225, 220, 210))
    tb = s4.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(2.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.RIGHT
    r_t = p_t.add_run()
    r_t.text = title
    r_t.font.name = 'Cairo'
    r_t.font.size = Pt(11)
    r_t.font.bold = True
    r_t.font.color.rgb = C_MUTED
    
    p_v = tf.add_paragraph()
    p_v.alignment = PP_ALIGN.RIGHT
    r_v = p_v.add_run()
    r_v.text = val
    r_v.font.name = 'Cairo'
    r_v.font.size = Pt(22)
    r_v.font.bold = True
    r_v.font.color.rgb = col
    
    p_d = tf.add_paragraph()
    p_d.alignment = PP_ALIGN.RIGHT
    r_d = p_d.add_run()
    r_d.text = desc
    r_d.font.name = 'Cairo'
    r_d.font.size = Pt(9.5)
    r_d.font.color.rgb = C_DARK

# =============================================================================
# SLIDE 5: REVENUE PERFORMANCE & SOURCES
# =============================================================================
s5 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s5, "مقارنة الإيرادات ومصادر الدخل (H1 2026 vs H1 2025)", "الأداء المالي")

# Table
tbl_shape = s5.shapes.add_table(7, 5, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
tbl = tbl_shape.table

# Column widths
tbl.columns[0].width = Inches(3.2)
tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(2.2)
tbl.columns[4].width = Inches(1.933)

rev_headers = ["بند الإيراد", "النصف الأول ٢٠٢٦م", "النصف الأول ٢٠٢٥م", "التغير (ريال)", "نسبة النمو"]
rev_rows = [
    ["أموال الزكاة المقيدة", "٧٠,٠٠٠ ريال", "٨٠,٠٠٠ ريال", "-١٠,٠٠٠", "-١٣٪"],
    ["علاج مقيد (مساعدات طبية)", "٧٥,٠٠٠ ريال", "٢٥,٠٠٠ ريال", "+٥٠,٠٠٠", "+٢٠٠٪"],
    ["المتجر الإلكتروني", "١٠,٤٦٩ ريال", "١٢٤ ريال", "+١٠,٣٤٥", "+٨,٣٤٣٪"],
    ["منصة تبرع الوطنية", "١,٢٠٣ ريال", "١٣,٧٨٦ ريال", "-١٢,٥٨٣", "-٩١٪"],
    ["تبرعات ودعم عام", "٤٠٧,٤٩٥ ريال", "٦٢,٥٦٤ ريال", "+٣٤٤,٩٣١", "+٥٥١٪"],
    ["اشتراكات العضوية", "١٨,٠٠٠ ريال", "١٨,٠٠٠ ريال", "٠", "٠٪"]
]

for c_idx, h in enumerate(rev_headers):
    cell = tbl.cell(0, c_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = h
    r.font.name = 'Cairo'
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

for r_idx, row in enumerate(rev_rows):
    bg = RGBColor(245, 242, 235) if r_idx % 2 == 1 else C_WHITE
    for c_idx, val in enumerate(row):
        cell = tbl.cell(r_idx + 1, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = 'Cairo'
        r.font.size = Pt(10.5)
        if c_idx == 4 and "+" in val:
            r.font.bold = True
            r.font.color.rgb = C_SUCCESS
        elif c_idx == 4 and "-" in val:
            r.font.bold = True
            r.font.color.rgb = C_DANGER

# =============================================================================
# SLIDE 6: BUDGET EXECUTION & PROGRESS
# =============================================================================
s6 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s6, "مستوى تنفيذ الموازنة التقديرية (٢٠٢٦م)", "التخطيط والالتزام المالي")

bud_data = [
    ("التبرعات والدعم (الإيرادات)", "٥٨٢,١٦٧ / ١,٥٢٧,٠٠٠ ريال", "٤٠.٠٢٪", C_SUCCESS),
    ("المساعدات العلاجية للمرضى", "٢٠٨,٦٠٥ / ٧٥٠,٠٠٠ ريال", "٢٧.٨١٪", C_SUCCESS),
    ("الرواتب والأجور والكادر", "١٤٤,٤٠٥ / ٤٧٢,٠٠٠ ريال", "٣٠.٥٩٪", C_PRIMARY),
    ("المصروفات التشغيلية", "١٠٩,٨٦٩ / ١٤٢,٣٠٠ ريال", "٧٧.٢١٪", C_WARNING),
    ("شراء الأصول والتجهيزات", "١٥,٦٢١ / ١٩,٤٥٠ ريال", "٨٠.٣١٪", C_DANGER)
]

for idx, (b_title, b_vals, b_pct, b_col) in enumerate(bud_data):
    top = 1.8 + idx * 0.95
    add_card(s6, 0.8, top, 11.733, 0.8, C_WHITE, RGBColor(225, 220, 210))
    tb = s6.shapes.add_textbox(Inches(1.0), Inches(top + 0.1), Inches(11.333), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    
    r_t = p.add_run()
    r_t.text = f"{b_title}:  "
    r_t.font.name = 'Cairo'
    r_t.font.size = Pt(11.5)
    r_t.font.bold = True
    r_t.font.color.rgb = C_PRIMARY
    
    r_v = p.add_run()
    r_v.text = f"{b_vals}  |  نسبة الإنجاز: "
    r_v.font.name = 'Cairo'
    r_v.font.size = Pt(11)
    r_v.font.color.rgb = C_MUTED
    
    r_p = p.add_run()
    r_p.text = f"{b_pct}"
    r_p.font.name = 'Cairo'
    r_p.font.size = Pt(13)
    r_p.font.bold = True
    r_p.font.color.rgb = b_col

# Summary callout
add_card(s6, 0.8, 6.4, 11.733, 0.7, C_PRIMARY, C_SECONDARY)
tb_sum = s6.shapes.add_textbox(Inches(1.0), Inches(6.45), Inches(11.333), Inches(0.6))
tf_sum = tb_sum.text_frame
p_sum = tf_sum.paragraphs[0]
p_sum.alignment = PP_ALIGN.CENTER
r_sum = p_sum.add_run()
r_sum.text = "إجمالي المنفذ الفعلي: ١,٠٦٠,٦٦٦ ريال من أصل موازنة سنوية ٢,٩٨١,٧٥٠ ريال (نسبة التنفيذ الإجمالية: ٣٥.٥٧٪)"
r_sum.font.name = 'Cairo'
r_sum.font.size = Pt(12)
r_sum.font.bold = True
r_sum.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 7: OPERATING EXPENSES STRUCTURE
# =============================================================================
s7 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s7, "هيكل المصروفات التشغيلية وأوجه الإنفاق", "الترشيد والشفافية")

# Left Column: Key Items
add_card(s7, 0.8, 1.8, 5.7, 5.0, C_WHITE, RGBColor(225, 220, 210))
tb_exp1 = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_exp1 = tb_exp1.text_frame
tf_exp1.word_wrap = True

p_expt1 = tf_exp1.paragraphs[0]
p_expt1.alignment = PP_ALIGN.RIGHT
r_expt1 = p_expt1.add_run()
r_expt1.text = "البنود التشغيلية الرئيسية (H1 2026):\n\n"
r_expt1.font.name = 'Cairo'
r_expt1.font.size = Pt(13)
r_expt1.font.bold = True
r_expt1.font.color.rgb = C_PRIMARY

items_exp1 = [
    "١. الرواتب الأساسية: ١٤٤,٤٠٥ ريال (٥٦.٨٪ من النفقات)",
    "٢. الإيجار المكتبي: ٦٣,٣٣٣ ريال (٢٤.٩٪ من النفقات)",
    "٣. التأمينات الاجتماعية: ١٤,٧٦٨ ريال (٥.٨٪)",
    "٤. أجور متعاونين ومحاسب: ١٣,٠٠٠ ريال (٥.١٪)",
    "٥. المحاسب القانوني للقوائم: ٤,٦٠٠ ريال (١.٨٪)",
    "٦. الكهرباء والخدمات: ٣,٨٦٧ ريال (١.٥٪)"
]
for it in items_exp1:
    p = tf_exp1.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = it + "\n"
    r.font.name = 'Cairo'
    r.font.size = Pt(10.5)
    r.font.color.rgb = C_DARK

# Right Column: Insights & Rent Savings
add_card(s7, 6.8, 1.8, 5.7, 5.0, C_PRIMARY_DARK, C_SECONDARY)
tb_exp2 = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_exp2 = tb_exp2.text_frame
tf_exp2.word_wrap = True

p_expt2 = tf_exp2.paragraphs[0]
p_expt2.alignment = PP_ALIGN.RIGHT
r_expt2 = p_expt2.add_run()
r_expt2.text = "ملاحظات الأداء والترشيد المؤسسي:\n\n"
r_expt2.font.name = 'Cairo'
r_expt2.font.size = Pt(13)
r_expt2.font.bold = True
r_expt2.font.color.rgb = C_SECONDARY

items_exp2 = [
    "• وفر الإيجار السنوي: الانتقال إلى مقر جديد بإيجار ٤٥ ألف بدلاً من ٧٠ ألف ريال (وفر ٢٥,٠٠٠ ريال سنوياً).",
    "• استقرار الالتزامات: سداد كامل الالتزامات المرحلة من ٢٠٢٥م بقيمة ١٨,٢١١ ريال.",
    "• استثمار التحول الرقمي: ٣,٠٠٠ ريال لتطوير البوابة الرسمية، وتفعيل نظام قيود السحابي.",
    "• التوصية الإدارية: ترشيد النفقات الإدارية لتكون متناسبة مع التوسع بالبرامج الطبية."
]
for it in items_exp2:
    p = tf_exp2.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = it + "\n"
    r.font.name = 'Cairo'
    r.font.size = Pt(10.5)
    r.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 8: MEDICAL PROGRAM "QUALITY OF LIFE" & 7 CASES
# =============================================================================
s8 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s8, "البرامج والخدمات الطبية للمرضى (برنامج جودة حياة)", "الأثر الميداني")

# 7 Patients Table
tbl_pat_shape = s8.shapes.add_table(8, 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tbl_p = tbl_pat_shape.table
tbl_p.columns[0].width = Inches(2.8)
tbl_p.columns[1].width = Inches(3.2)
tbl_p.columns[2].width = Inches(3.5)
tbl_p.columns[3].width = Inches(2.233)

pat_hdrs = ["اسم المستفيد", "الجهة العلاجية", "التشخيص الطبي", "مبلغ الدعم (ريال)"]
pat_rows = [
    ["فايز أحمد عبدالعزيز", "المستشفى السعودي الألماني", "سرطان الدم (علاج مناعي وكيماوي)", "١٥٠,٠٠٠ ريال"],
    ["زينب عمر علي", "المستشفى السعودي الألماني", "سرطان نخر العظم", "٣٠,٠٠٠ ريال"],
    ["كندفة محمد عتبة", "مدينة الملك سلمان الطبية", "تنويم ورعاية تحت الملاحظة", "٧,٠٠٠ ريال"],
    ["شوق حسن الأنور", "المستشفى السعودي الألماني", "منظار جراحي متقدم", "٧,٠٠٠ ريال"],
    ["سامية سليمان محمد", "مستشفى المواساة بالمدينة", "استئصال كتلة ورمية بالصدر", "٦,٣٥٠ ريال"],
    ["زبيدة شمس الدين خاتم", "المستشفى السعودي الألماني", "ورم بالقولون", "٦,٣٣٠.٣١ ريال"],
    ["محمد أحمد الشرفي", "مستشفى المواساة بالمدينة", "أشعة رنين مغناطيسي", "١,٩٢٥ ريال"]
]

for c_idx, h in enumerate(pat_hdrs):
    cell = tbl_p.cell(0, c_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = h
    r.font.name = 'Cairo'
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

for r_idx, row in enumerate(pat_rows):
    bg = RGBColor(245, 242, 235) if r_idx % 2 == 1 else C_WHITE
    for c_idx, val in enumerate(row):
        cell = tbl_p.cell(r_idx + 1, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = 'Cairo'
        r.font.size = Pt(10)
        if c_idx == 3:
            r.font.bold = True
            r.font.color.rgb = C_PRIMARY

# =============================================================================
# SLIDE 9: CASE REJECTION ANALYSIS (14 CASES)
# =============================================================================
s9 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s9, "تحليل الحالات غير المقبولة (١٤ حالة) والتوصيات", "حوكمة المساعدات")

rej_boxes = [
    ("انتهاء صلاحية الإقامة (٧ حالات - ٥٠٪)", "بسمة هارون، سيد الأمين، فريدة عظيم، عطور عباس، هاجر الصادق، عبدالله دياب، أحمد خير\n• التوصية: التنسيق مع المانحين لقبول الحالات الإنسانية الحرجة وتعديل اللائحة."),
    ("تغطية كاملة من جمعية أخرى (حالتان - ١٤.٣٪)", "هديباء الجهني (مياه بيضاء)، علي قايد (شرايين)\n• التوصية: تفعيل الربط الإلكتروني لتفادي الازدواجية وسرعة توجيه الدعم لمرضى آخرين."),
    ("أخطاء بالتقرير الطبي والتواريخ (حالتان - ١٤.٣٪)", "ريم فواز (ورم ليفي)، جوهرة خان (أخطاء تواريخ)\n• التوصية: توجيه المستفيدين لتصحيح المستندات الطبية وإعادة الرفع."),
    ("تأمين طبي ساري / انتهاء تأشيرة (٣ حالات - ٢١.٤٪)", "فؤاد لطف (تأمين طبي)، حمزة هندية (انتهاء زيارة)، مزاهر الهادي (قيد تسليم التعميد).")
]

for idx, (rtitle, rdesc) in enumerate(rej_boxes):
    r_i = idx // 2
    c_i = idx % 2
    left = 0.8 + c_i * 5.95
    top = 1.8 + r_i * 2.5
    
    add_card(s9, left, top, 5.7, 2.3, C_WHITE, RGBColor(225, 220, 210))
    tb = s9.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(5.4), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.RIGHT
    r_t = p_t.add_run()
    r_t.text = rtitle + "\n"
    r_t.font.name = 'Cairo'
    r_t.font.size = Pt(11.5)
    r_t.font.bold = True
    r_t.font.color.rgb = C_DANGER if idx == 0 else C_PRIMARY
    
    p_d = tf.add_paragraph()
    p_d.alignment = PP_ALIGN.RIGHT
    r_d = p_d.add_run()
    r_d.text = rdesc
    r_d.font.name = 'Cairo'
    r_d.font.size = Pt(10)
    r_d.font.color.rgb = C_DARK

# =============================================================================
# SLIDE 10: 9 HEALTHCARE PARTNERSHIPS
# =============================================================================
s10 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s10, "شبكة الشراكات الصحية والمؤسسية (٩ جهات)", "التكامل والشراكات")

partners = [
    ("المستشفى السعودي الألماني", "علاج الأورام وسرطانات الدم والمناظير التخصصية"),
    ("مستشفى المواساة بالمدينة", "العمليات الجراحية المتقدمة والأشعة والرنين"),
    ("مدينة الملك سلمان الطبية", "الرعاية المرجعية التخصصية والتنويم"),
    ("مستشفى د. حامد الأحمدي", "جراحات اليوم الواحد والعيادات الاستشارية"),
    ("مستشفى المدينة الوطني", "خدمات الطوارئ والملاحظة والتحاليل"),
    ("مستشفى المدينة الطبي العام", "الفحوصات العامة ورعاية الأمراض المزمنة"),
    ("مستشفى واد الطبي", "علاج الإصابات الرياضية وجراحة العظام"),
    ("شركة مداواة ورعاية الطبية", "توفير الأدوية والمستلزمات الطبية المنزلية"),
    ("جمعية جَنَى لتأهيل المعاقات", "التأهيل الطبي والتكامل مع ذوي الإعاقة")
]

for idx, (pname, pscope) in enumerate(partners):
    r_i = idx // 3
    c_i = idx % 3
    left = 0.8 + c_i * 3.95
    top = 1.8 + r_i * 1.65
    
    add_card(s10, left, top, 3.8, 1.5, C_WHITE, RGBColor(225, 220, 210))
    tb = s10.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(3.5), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_n = tf.paragraphs[0]
    p_n.alignment = PP_ALIGN.RIGHT
    r_n = p_n.add_run()
    r_n.text = pname + "\n"
    r_n.font.name = 'Cairo'
    r_n.font.size = Pt(11)
    r_n.font.bold = True
    r_n.font.color.rgb = C_PRIMARY
    
    p_s = tf.add_paragraph()
    p_s.alignment = PP_ALIGN.RIGHT
    r_s = p_s.add_run()
    r_s.text = pscope
    r_s.font.name = 'Cairo'
    r_s.font.size = Pt(9.5)
    r_s.font.color.rgb = C_MUTED

# =============================================================================
# SLIDE 11: HR, 100% SAUDIZATION & TRAINING
# =============================================================================
s11 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s11, "الموارد البشرية والنمو المؤسسي", "رأس المال البشري")

# Left Column: Staff
add_card(s11, 0.8, 1.8, 5.7, 5.0, C_WHITE, RGBColor(225, 220, 210))
tb_hr1 = s11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_hr1 = tb_hr1.text_frame
tf_hr1.word_wrap = True

p_hrt1 = tf_hr1.paragraphs[0]
p_hrt1.alignment = PP_ALIGN.RIGHT
r_hrt1 = p_hrt1.add_run()
r_hrt1.text = "الهيكل الوظيفي والكادر التنفيذي:\n\n"
r_hrt1.font.name = 'Cairo'
r_hrt1.font.size = Pt(13)
r_hrt1.font.bold = True
r_hrt1.font.color.rgb = C_PRIMARY

items_hr1 = [
    "• أ. بيان سعد المحمدي: المدير التنفيذي (موظف).",
    "• أ. غدير أحمد الحربي: المسؤول المالي والمشرفة على البرامج (موظفة).",
    "• أ. طراد محمد سمان: سكرتير تنفيذي (موظف).",
    "• أ. محمد الحسن بشير: محاسب متعاون.",
    "• أ. فيصل الجهني: مسؤول علاقات وإعلام (متعاون قيد الترسيم).",
    "• نسبة التوطين: ١٠٠٪ في كافة المناصب الرسمية."
]
for it in items_hr1:
    p = tf_hr1.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = it + "\n"
    r.font.name = 'Cairo'
    r.font.size = Pt(10.5)
    r.font.color.rgb = C_DARK

# Right Column: Training & Volunteering
add_card(s11, 6.8, 1.8, 5.7, 5.0, C_PRIMARY, C_SECONDARY)
tb_hr2 = s11.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
tf_hr2 = tb_hr2.text_frame
tf_hr2.word_wrap = True

p_hrt2 = tf_hr2.paragraphs[0]
p_hrt2.alignment = PP_ALIGN.RIGHT
r_hrt2 = p_hrt2.add_run()
r_hrt2.text = "مؤشرات التأهيل والتدريب المؤسسي:\n\n"
r_hrt2.font.name = 'Cairo'
r_hrt2.font.size = Pt(13)
r_hrt2.font.bold = True
r_hrt2.font.color.rgb = C_SECONDARY

items_hr2 = [
    "• ٨ دورات تدريبية متخصصة: استفاد منها موظفان بالكادر لرفع الكفاءة المالية والإدارية (مقابل ٠ دورات في ٢٠٢٥).",
    "• أنظمة الموارد البشرية: تطبيق نظام البصمة الإلكترونية وضبط الدوام والأرشفة الرقمية لكافة ملفات الموظفين.",
    "• ملف التطوع: تسجيل ٤ فرص تطوعية بالنصف الأول، مع التوصية بإعادة تفعيل الفرق التطوعية لرفع العائد المجتمعي."
]
for it in items_hr2:
    p = tf_hr2.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = it + "\n"
    r.font.name = 'Cairo'
    r.font.size = Pt(10.5)
    r.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 12: RESOURCE DEVELOPMENT & 27 GRANTS
# =============================================================================
s12 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s12, "تنمية الموارد ومسار طلبات المنح (٢٧ جهة)", "الاستدامة والتمويل")

grants_summary = [
    ("منح مقبولة ومحققة (٤٠,٠٠٠ ريال)", "• مؤسسة إبراهيم العنقري: ٢٠,٠٠٠ ريال.\n• وقف الشيخ عبدالعزيز أبو زيد: ٢٠,٠٠٠ ريال.\n(تحققتا رغم عدم اكتمال الحوكمة والقوائم حينها).", C_SUCCESS),
    ("طلبات قيد الدراسة والمتابعة (١١ جهة)", "• صندوق دعم الجمعيات (مشروع جودة حياة).\n• أوقاف الشيخ صالح الراجحي (دفء وغذاء ٢٠٢٧).\n• بنك البلاد، بنك الرياض، و ٦ شركات لحجاج الداخل.", C_WARNING),
    ("اعتذارات بسبب الموازنة أو الشروط (٨ جهات)", "• انتهاء الموازنات: مؤسسة الماجد، مؤسسة الشاوي، شركة طيبة للاستثمار، مجموعة فنادق.\n• اشتراط الحوكمة: أوقاف الضحيان، مؤسسة طلال، مؤسسة الحمدان، مؤسسة المهيدب.", C_DANGER)
]

for idx, (gtitle, gdesc, gcol) in enumerate(grants_summary):
    left = 0.8 + idx * 3.95
    add_card(s12, left, 1.8, 3.8, 5.0, C_WHITE, RGBColor(225, 220, 210))
    tb = s12.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(3.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.RIGHT
    r_t = p_t.add_run()
    r_t.text = gtitle + "\n\n"
    r_t.font.name = 'Cairo'
    r_t.font.size = Pt(12)
    r_t.font.bold = True
    r_t.font.color.rgb = gcol
    
    p_d = tf.add_paragraph()
    p_d.alignment = PP_ALIGN.RIGHT
    r_d = p_d.add_run()
    r_d.text = gdesc
    r_d.font.name = 'Cairo'
    r_d.font.size = Pt(10)
    r_d.font.color.rgb = C_DARK

# =============================================================================
# SLIDE 13: GOVERNANCE & 3-PHASE ROADMAP
# =============================================================================
s13 = add_blank_slide(RGBColor(248, 247, 244))
add_header(s13, "الحوكمة وخارطة طريق النصف الثاني (٣ مراحل)", "التحول المؤسسي")

phases = [
    ("المرحلة الأولى (الشهر الأول)\nاستكمال الحوكمة ومنصة نوى", "• استيفاء متطلبات معيار الامتثال والحوكمة.\n• رفع درجة تقييم الجمعية المعتمدة.\n• تفعيل منصة نوى للمنح والشراكات.\n• توظيف القوائم المالية المعتمدة."),
    ("المرحلة الثانية (الشهر الثاني)\nتنمية الموارد وبطاقة طبيبي", "• إطلاق مبادرة «بطاقة طبيبي» للمزايا الطبية.\n• تعديل لائحة المساعدات لرفع نسبة القبول.\n• بناء قاعدة بيانات المانحين والأوقاف.\n• استقطاب كفاءات تنمية الموارد."),
    ("المرحلة الثالثة (الشهر الثالث)\nالاستعداد المبكر لـ Q1 2027", "• تجهيز الحقائب الاستثمارية للصناديق الكبرى.\n• استهداف المانحين الذين أغلقت موازناتهم.\n• تفعيل دور الجمعية العمومية والمجلس.\n• ربط المؤشرات بالأثر الصحي والمجتمعي.")
]

for idx, (p_title, p_desc) in enumerate(phases):
    left = 0.8 + idx * 3.95
    add_card(s13, left, 1.8, 3.8, 5.0, C_PRIMARY_DARK, C_SECONDARY)
    tb = s13.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(3.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.RIGHT
    r_t = p_t.add_run()
    r_t.text = p_title + "\n\n"
    r_t.font.name = 'Cairo'
    r_t.font.size = Pt(12)
    r_t.font.bold = True
    r_t.font.color.rgb = C_SECONDARY
    
    p_d = tf.add_paragraph()
    p_d.alignment = PP_ALIGN.RIGHT
    r_d = p_d.add_run()
    r_d.text = p_desc
    r_d.font.name = 'Cairo'
    r_d.font.size = Pt(10)
    r_d.font.color.rgb = C_WHITE

# =============================================================================
# SLIDE 14: CLOSING & OFFICIAL CONTACTS
# =============================================================================
s14 = add_blank_slide(C_PRIMARY_DARK)

tb_end = s14.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.8))
tf_e = tb_end.text_frame
tf_e.word_wrap = True

pe1 = tf_e.paragraphs[0]
pe1.alignment = PP_ALIGN.CENTER
re1 = pe1.add_run()
re1.text = "شكراً لثقتكم ودعمكم المستمر لرسالة جمعية طبيبي الأهلية\n\n"
re1.font.name = 'Cairo'
re1.font.size = Pt(26)
re1.font.bold = True
re1.font.color.rgb = C_SECONDARY

pe2 = tf_e.add_paragraph()
pe2.alignment = PP_ALIGN.CENTER
re2 = pe2.add_run()
re2.text = "نسعى لمواصلة المسيرة نحو قطاع صحي غير ربحي مستدام ومؤثر في طيبة الطيبة\n\n"
re2.font.name = 'Cairo'
re2.font.size = Pt(14)
re2.font.color.rgb = C_WHITE

pe3 = tf_e.add_paragraph()
pe3.alignment = PP_ALIGN.CENTER
re3 = pe3.add_run()
re3.text = "هاتف: 00966555606347  |  البريد: tabibi2025med@gmail.com  |  المدينة المنورة - حي الفتح\nالمركز الوطني لتنمية القطاع غير الربحي - ترخيص: ١٠٠٠٧٣٠٧٠٠"
re3.font.name = 'Cairo'
re3.font.size = Pt(12)
re3.font.color.rgb = C_SECONDARY_LT

prs.save(output_pptx)
print(f"Generated official PowerPoint presentation successfully: {output_pptx}")

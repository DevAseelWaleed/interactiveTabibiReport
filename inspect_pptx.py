# -*- coding: utf-8 -*-
import os, sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

base_dir = os.path.dirname(os.path.abspath(__file__))

pptx_files = [f for f in os.listdir(base_dir) if f.endswith('.pptx')]

for pf in pptx_files:
    full_path = os.path.join(base_dir, pf)
    print(f"\n{'='*70}\nREADING PPTX: {pf}\n{'='*70}")
    try:
        prs = Presentation(full_path)
        print(f"Total slides: {len(prs.slides)}")
        for idx, slide in enumerate(prs.slides):
            print(f"\n--- Slide {idx+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            print(f"  [Text] {text}")
                elif shape.has_table:
                    print("  [Table]:")
                    for row in shape.table.rows:
                        row_txt = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        print(f"    {row_txt}")
    except Exception as e:
        print(f"Error reading {pf}: {e}")

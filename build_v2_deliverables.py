# -*- coding: utf-8 -*-
import os, sys, shutil
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import parse_xml, OxmlElement
from docx.oxml.ns import nsdecls, qn

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')

base_dir = r"e:\Work\زبون تقرير نصف سنوي طبيبي"
v2_dir = os.path.join(base_dir, "التقرير_الاحترافي_المطور")
v1_dir = os.path.join(base_dir, "التقرير_الجديد")

PRIMARY = RGBColor(107, 29, 58)     # #6B1D3A
SECONDARY = RGBColor(201, 169, 110) # #C9A96E
DARK_BG = RGBColor(46, 11, 23)
TEXT_DARK = RGBColor(45, 45, 45)

def set_cell_background(cell, fill_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
    tcPr.append(shd)

def add_rtl_heading(doc, text, level):
    p = doc.add_heading(text, level=level)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p.paragraph_format.bidi = True
    for r in p.runs:
        r.font.name = 'Arial'
        if level == 1:
            r.font.color.rgb = PRIMARY
            r.font.size = Pt(16)
            r.bold = True
        elif level == 2:
            r.font.color.rgb = PRIMARY
            r.font.size = Pt(13)
            r.bold = True
    return p

def add_rtl_p(doc, text, bold=False, color=None, size=11, align=WD_ALIGN_PARAGRAPH.RIGHT, italic=False):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.bidi = True
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = 1.2
    r = p.add_run(text)
    r.font.name = 'Arial'
    r.bold = bold
    r.italic = italic
    r.font.size = Pt(size)
    if color:
        r.font.color.rgb = color
    else:
        r.font.color.rgb = TEXT_DARK
    return p

def build_word_document(docx_path):
    print(f"Building Word: {docx_path}")
    doc = Document()
    for s in doc.sections:
        s.page_width = Inches(8.27)
        s.page_height = Inches(11.69)
        s.top_margin = Inches(0.8)
        s.bottom_margin = Inches(0.8)
        s.left_margin = Inches(0.8)
        s.right_margin = Inches(0.8)

    # 1. Cover
    add_rtl_p(doc, "جمعية طبيبي الأهلية بالمدينة المنورة", bold=True, color=SECONDARY, size=15, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_rtl_p(doc, "التقرير النصف سنوي الشامل لعام ٢٠٢٦م", bold=True, color=PRIMARY, size=24, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_rtl_p(doc, "«مدعم بالتدقيق الأدائي ومطابقة مستهدفات الخطة الاستراتيجية والتشغيلية»", bold=True, color=PRIMARY, size=13, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_rtl_p(doc, "الفترة من ١ يناير إلى ٣٠ يونيو ٢٠٢٦م | ترخيص رقم: (١٠٠٠٧٣٠٧٠٠)", color=TEXT_DARK, size=11, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_rtl_p(doc, "« ثـقـة • أثــر • اسـتـدامـة »", bold=True, color=SECONDARY, size=14, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()

    # 2. Table of Contents
    add_rtl_heading(doc, "فهرس المحتويات التنفيذي", level=1)
    toc_items = [
        ("١. المقدمة والقيادة الرشيدة ومجلس الإدارة", "ص ٢"),
        ("٢. كلمة رئيس مجلس الإدارة", "ص ٣"),
        ("٣. الملخص التنفيذي ومؤشرات الأداء الرئيسية", "ص ٤"),
        ("٤. التدقيق الأدائي والتقييم الاستراتيجي الشامل (BSC)", "ص ٥"),
        ("٥. مصفوفة مطابقة الخطة الاستراتيجية بالمنجز الفعلي (١٤ مؤشراً)", "ص ٦"),
        ("٦. حالة محفظة البرامج الاستراتيجية الـ (١٠)", "ص ٧"),
        ("٧. تحليل الفجوات الاستراتيجية وإدارة المخاطر", "ص ٨"),
        ("٨. التحليل والأداء المالي والموازنة التشغيلية", "ص ٩"),
        ("٩. البرامج والمساعدات الطبية وخدمة المستفيدين", "ص ١٠"),
        ("١٠. منظومة الشراكات الصحية والمؤسسية (٩ شركاء)", "ص ١١"),
        ("١١. الموارد البشرية والتحول المؤسسي والحوكمة", "ص ١٢"),
        ("١٢. قصص الأثر الإنساني ورسائل المستفيدين", "ص ١٣"),
        ("١٣. خارطة طريق التصحيح والتوصيات الاستراتيجية (H2 2026)", "ص ١٤"),
        ("١٤. الملاحق التفصيلية (الداعمون، والأصول الثابتة)", "ص ١٥")
    ]
    t_toc = doc.add_table(rows=len(toc_items)+1, cols=2)
    t_toc.alignment = WD_TABLE_ALIGNMENT.CENTER
    t_toc.rows[0].cells[0].paragraphs[0].text = "البند / المحور"
    t_toc.rows[0].cells[1].paragraphs[0].text = "رقم الصفحة"
    for c in t_toc.rows[0].cells:
        set_cell_background(c, "6B1D3A")
        c.paragraphs[0].runs[0].font.color.rgb = RGBColor(255,255,255)
        c.paragraphs[0].runs[0].font.bold = True
        c.paragraphs[0].paragraph_format.bidi = True
    for i, (item, page) in enumerate(toc_items):
        r = t_toc.rows[i+1]
        r.cells[0].paragraphs[0].text = item
        r.cells[1].paragraphs[0].text = page
        for c in r.cells:
            c.paragraphs[0].paragraph_format.bidi = True
            if i % 2 == 1:
                set_cell_background(c, "F8F6F0")
    doc.add_page_break()

    # 3. Royal Leadership
    add_rtl_heading(doc, "القيادة الرشيدة ومجلس الإدارة", level=1)
    leaders = [
        ("صاحب السمو الملكي الأمير محمد بن سلمان بن عبدالعزيز", "ولي العهد رئيس مجلس الوزراء", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة والثقافة والمجالات البحثية، وسنعتمد عليه بشكل رئيسي.»"),
        ("خادم الحرمين الشريفين الملك سلمان بن عبدالعزيز آل سعود", "ملك المملكة العربية السعودية", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، وما نراه من مؤسسات خيرية في مختلف المجالات… إلا جانبًا من الجوانب المشرقة لبلادنا.»"),
        ("صاحب السمو الملكي الأمير سلمان بن سلطان بن عبدالعزيز", "أمير منطقة المدينة المنورة", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية على مستوى المنطقة باعتبارها شريكًا استراتيجيًا للقطاعين العام والخاص في تحسين جودة الحياة.»")
    ]
    for name, title, quote in leaders:
        add_rtl_p(doc, f"• {name} - {title}", bold=True, color=PRIMARY, size=12)
        add_rtl_p(doc, quote, italic=True, color=TEXT_DARK, size=10.5)
    doc.add_page_break()

    # 4. Strategic Audit & BSC Score
    add_rtl_heading(doc, "التدقيق الأدائي والتقييم الاستراتيجي الشامل (BSC Evaluation)", level=1)
    add_rtl_p(doc, "بناءً على الفحص المستندي المقارن بين مستهدفات الخطة الاستراتيجية والتشغيلية والمنجزات الفعلية للنصف الأول ٢٠٢٦م:", size=11)
    
    t_bsc = doc.add_table(rows=6, cols=4)
    t_bsc.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers_bsc = ["المحور الاستراتيجي (BSC)", "الوزن النسبي", "نسبة إنجاز H1", "الدرجة الموزونة المحققة"]
    for j, h in enumerate(headers_bsc):
        t_bsc.rows[0].cells[j].paragraphs[0].text = h
        set_cell_background(t_bsc.rows[0].cells[j], "6B1D3A")
        t_bsc.rows[0].cells[j].paragraphs[0].runs[0].font.color.rgb = RGBColor(255,255,255)
        t_bsc.rows[0].cells[j].paragraphs[0].runs[0].font.bold = True
        t_bsc.rows[0].cells[j].paragraphs[0].paragraph_format.bidi = True
    
    bsc_data = [
        ("١. محور الأثر والبرامج الطبية", "٤٠٪", "١٣.٩٥٪", "٥.٥٨ من ٤٠"),
        ("٢. المحور المالي والموازنة التشغيلية", "٣٠٪", "٣٢.٩٦٪", "٩.٨٩ من ٣٠"),
        ("٣. محور الشراكات والعمليات الداخلية", "١٥٪", "٥٥.٠٠٪", "٨.٢٥ من ١٥"),
        ("٤. محور الحوكمة والمؤسسية", "١٥٪", "٦٠.٠٠٪", "٩.٠٠ من ١٥"),
        ("الإجمالي العام الموزون للأداء الاستراتيجي (H1 2026)", "١٠٠٪", "٣٢.٧٢٪", "٣٢.٧٢ من ١٠٠ (يحتاج تحسين)")
    ]
    for i, row in enumerate(bsc_data):
        r = t_bsc.rows[i+1]
        for j, val in enumerate(row):
            r.cells[j].paragraphs[0].text = val
            r.cells[j].paragraphs[0].paragraph_format.bidi = True
            if i == 4:
                set_cell_background(r.cells[j], "F4EFE6")
                r.cells[j].paragraphs[0].runs[0].font.bold = True
                r.cells[j].paragraphs[0].runs[0].font.color.rgb = PRIMARY
            elif i % 2 == 1:
                set_cell_background(r.cells[j], "F8F6F0")

    add_rtl_p(doc, "التقييم العام المعتمد: «يحتاج إلى تحسين جذري وإعادة ضبط مسار (Needs Significant Improvement & Realignment)» نظراً لفجوة المستفيدين (٧ مرضى فقط مقابل مستهدف ٣٦ ألف) وتعطيل ٩ برامج معتمدة.", bold=True, color=PRIMARY, size=11)
    doc.add_page_break()

    # 5. Plan vs Actual Matrix
    add_rtl_heading(doc, "مصفوفة مطابقة الخطة الاستراتيجية بالمنجز الفعلي (١٤ مؤشراً معتمداً)", level=1)
    matrix_headers = ["م", "الهدف / النشاط", "المؤشر المعتمد", "المستهدف السنوي", "مستهدف H1", "المنجز الفعلي", "نسبة الإنجاز", "الحالة"]
    t_mat = doc.add_table(rows=15, cols=8)
    t_mat.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, h in enumerate(matrix_headers):
        t_mat.rows[0].cells[j].paragraphs[0].text = h
        set_cell_background(t_mat.rows[0].cells[j], "6B1D3A")
        t_mat.rows[0].cells[j].paragraphs[0].runs[0].font.color.rgb = RGBColor(255,255,255)
        t_mat.rows[0].cells[j].paragraphs[0].runs[0].font.bold = True
        t_mat.rows[0].cells[j].paragraphs[0].paragraph_format.bidi = True
        t_mat.rows[0].cells[j].paragraphs[0].runs[0].font.size = Pt(9)
    
    matrix_rows = [
        ("١", "الإيرادات الكلية للجمعية", "إجمالي الدخل (ريال)", "٦,٨٤٦,٠٠٠", "٣,٤٢٣,٠٠٠", "٥٨٢,١٦٧.٥٢", "١٧.٠١٪", "متأخر حرِج"),
        ("٢", "إيرادات الموازنة التشغيلية", "إيراد الموازنة (ريال)", "١,٥٢٧,٠٠٠", "١,٥٢٧,٠٠٠", "٥٨٢,١٦٧.٥٢", "٣٨.١٢٪", "متأخر"),
        ("٣", "المساعدات العلاجية المباشرة", "مبالغ المساعدات (ريال)", "١,٥٠٠,٠٠٠", "٧٥٠,٠٠٠", "٢٠٨,٦٠٥.٠٠", "٢٧.٨١٪", "متأخر حرِج"),
        ("٤", "المستفيدون من الخدمات الصحية", "عدد المستفيدين (فرد)", "٣٦,٦٠٦", "١٨,٣٠٣", "٧ مستفيدين", "٠.٠٣٨٪", "متعثر تماماً"),
        ("٥", "الاستشارات الطبية والدوائية", "عدد الاستشارات", "١,٢٠٠", "٦٠٠", "٠ استشارة", "٠.٠٠٪", "لم يبدأ"),
        ("٦", "الدراسات واستطلاعات الرأي", "عدد الدراسات", "٦ دراسات", "٣ دراسات", "٠ دراسة", "٠.٠٠٪", "لم يبدأ"),
        ("٧", "ساعات وقيمة التطوع", "ساعات وقيمة التطوع", "٣,٠٠٠ س (٢٠٢ ألف)", "١,٥٠٠ س (١٠١ ألف)", "٤ فرص تطوعية", "غير مدققة", "متعثر"),
        ("٨", "تفعيل محفظة البرامج", "عدد البرامج النشطة", "١٠ برامج", "١٠ برامج", "برنامج واحد فقط", "١٠.٠٠٪", "متعثر"),
        ("٩", "عقد الشراكات الصحية", "عدد الشراكات", "٩ شراكات", "٩ شراكات", "٩ شراكات مفعلة", "١٠٠.٠٠٪", "مكتمل"),
        ("١٠", "توطين الوظائف والكادر", "نسبة التوطين (٪)", "١٠٠٪", "١٠٠٪", "١٠٠٪ (٣ موظفين)", "١٠٠.٠٠٪", "مكتمل"),
        ("١١", "تدريب وتأهيل الكادر", "عدد الدورات", "٤ دورات", "٢ دورة", "٨ دورات", "٤٠٠.٠٠٪", "متقدم ومكتمل"),
        ("١٢", "التحول الرقمي والمحاسبي", "تطبيق نظام سحابي", "نظام قيود", "نظام قيود", "تم تشغيل قيود", "١٠٠.٠٠٪", "مكتمل"),
        ("١٣", "معايير الحوكمة ومنصة نوى", "نسبة الامتثال", "١٠٠٪", "٥٠٪", "طلب استشاري", "٢٥.٠٠٪", "متأخر"),
        ("١٤", "تنويع مصادر الدخل الذاتي", "عدد مصادر الدخل", "٦ مصادر", "٦ مصادر", "٦ مصادر نشطة", "٦٦.٦٧٪", "قيد التنفيذ")
    ]
    for i, row in enumerate(matrix_rows):
        r = t_mat.rows[i+1]
        for j, val in enumerate(row):
            r.cells[j].paragraphs[0].text = val
            r.cells[j].paragraphs[0].paragraph_format.bidi = True
            r.cells[j].paragraphs[0].runs[0].font.size = Pt(8.5)
            if i % 2 == 1:
                set_cell_background(r.cells[j], "F8F6F0")
    doc.add_page_break()

    # 6. Critical Gaps & Roadmap
    add_rtl_heading(doc, "تحليل الفجوات الاستراتيجية والتوصيات التصحيحية (H2 2026)", level=1)
    gaps = [
        ("الفجوة الأولى: انحسار نطاق المستفيدين", "خدمة ٧ مرضى فقط من أصل ٣٦,٦٠٦ مستهدفين نتيجة حصر الصرف في الجراحات المعقدة وإيقاف العيادات الوقائية والاستشارات."),
        ("الفجوة الثانية: تشدد لائحة المساعدات العلاجية", "رفض ٦٦.٧٪ من الحالات المتقدمة (١٤ حالة) لانتهاء الإقامة أو نقص التقارير مما عطل صرف موازنة العلاج."),
        ("الفجوة الثالثة: تركز الإيرادات وضعف المنح", "اعتماد الجمعية بنسبة ٤٣٪ على تبرع فردي واحد (٢٥٠ ألف ريال) وتدني نسبة نجاح طلبات المنح إلى ٧.٤٪."),
        ("الفجوة الرابعة: تراجع النشاط التطوعي", "تنفيذ ٤ فرص تطوعية فقط دون توثيق للساعات مقابل مستهدف ٣,٠٠٠ ساعة بقيمة ٢٠٢.٥ ألف ريال.")
    ]
    for title, desc in gaps:
        add_rtl_p(doc, f"• {title}:", bold=True, color=PRIMARY, size=11.5)
        add_rtl_p(doc, desc, color=TEXT_DARK, size=10.5)

    add_rtl_heading(doc, "خارطة طريق التوصيات للنصف الثاني ٢٠٢٦م", level=2)
    recs = [
        ("١. تعديل عاجل للائحة المساعدات", "وضع استثناءات إنسانية لرفع نسبة القبول وصرف المتبقي من الموازنة (٥٤١ ألف ريال)."),
        ("٢. تشغيل برامج الاستشارات والفحوصات", "إطلاق الاستشارات الهاتفية والقوافل الميدانية لخدمة آلاف المستفيدين بتكلفة منخفضة."),
        ("٣. اعتماد موازنة الحوكمة ونوى", "التعاقد مع الفريق الاستشاري (١٥-٢١ ألف ريال) لإنهاء متطلبات الامتثال وفتح المنح الكبرى."),
        ("٤. إطلاق «بطاقة طبيبي» وتوحيد الخطة", "استثمار شبكة المستشفيات لإصدار بطاقة مزايا صحية وتوحيد مستهدفات الخطة الاستراتيجية.")
    ]
    for title, desc in recs:
        add_rtl_p(doc, f"• {title}:", bold=True, color=PRIMARY, size=11.5)
        add_rtl_p(doc, desc, color=TEXT_DARK, size=10.5)

    doc.save(docx_path)
    print(f"Saved Word Document: {docx_path}")

def build_powerpoint(pptx_path):
    print(f"Building PPTX: {pptx_path}")
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]

    C_PRIMARY = PRGBColor(107, 29, 58)
    C_SECONDARY = PRGBColor(201, 169, 110)
    C_WHITE = PRGBColor(255, 255, 255)
    C_TEXT = PRGBColor(45, 45, 45)

    def add_header(slide, title, subtitle):
        tb = slide.shapes.add_textbox(PInches(0.8), PInches(0.4), PInches(11.733), PInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(22)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = PPt(12)
        p2.font.color.rgb = C_SECONDARY
        p2.alignment = PP_ALIGN.RIGHT

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank_layout)
    tb = s1.shapes.add_textbox(PInches(1.0), PInches(1.5), PInches(11.333), PInches(4.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "جمعية طبيبي الأهلية بالمدينة المنورة"
    p.font.size = PPt(20)
    p.font.color.rgb = C_SECONDARY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "التقرير النصف سنوي الشامل لعام ٢٠٢٦م"
    p2.font.size = PPt(32)
    p2.font.bold = True
    p2.font.color.rgb = C_PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "«مدعم بالتدقيق الأدائي ومطابقة مستهدفات الخطة الاستراتيجية والتشغيلية»"
    p3.font.size = PPt(16)
    p3.font.bold = True
    p3.font.color.rgb = C_PRIMARY
    p3.alignment = PP_ALIGN.CENTER

    # Slide 2: Royal Leadership
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "القيادة الرشيدة ومجلس الإدارة", "الالتزام بتوجيهات القيادة في تمكين القطاع الصحي غير الربحي")
    leaders_info = [
        ("صاحب السمو الملكي الأمير محمد بن سلمان", "ولي العهد رئيس مجلس الوزراء", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة.»", 0.8),
        ("خادم الحرمين الشريفين الملك سلمان بن عبدالعزيز", "ملك المملكة العربية السعودية", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه ومؤسساتها الخيرية.»", 4.8),
        ("صاحب السمو الملكي الأمير سلمان بن سلطان", "أمير منطقة المدينة المنورة", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية كشريك استراتيجي في جودة الحياة.»", 8.8)
    ]
    for name, title, quote, left in leaders_info:
        tb = s2.shapes.add_textbox(PInches(left), PInches(1.6), PInches(3.7), PInches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_SECONDARY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = f"\n{quote}"
        p3.font.size = PPt(11)
        p3.font.italic = True
        p3.font.color.rgb = C_TEXT
        p3.alignment = PP_ALIGN.RIGHT

    # Slide 3: Executive Strategic Audit Summary (BSC)
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "التدقيق الأدائي والتقييم الاستراتيجي الشامل (BSC)", "نتائج تقييم النصف الأول ٢٠٢٦م وفق مناظير بطاقة الأداء المتوازن")
    bsc_cards = [
        ("محور الأثر والبرامج الطبية (٤٠٪)", "١٣.٩٥٪", "٥.٥٨ من ٤٠ نقطة | خدمة ٧ مرضى من أصل ٣٦ ألف", 0.8),
        ("المحور المالي والموازنة (٣٠٪)", "٣٢.٩٦٪", "٩.٨٩ من ٣٠ نقطة | تحصيل ٥٨٢ ألف ر.س (+١٩٢٪)", 3.8),
        ("محور الشراكات والعمليات (١٥٪)", "٥٥.٠٠٪", "٨.٢٥ من ١٥ نقطة | ٩ شراكات مستشفيات فاعلة", 6.8),
        ("محور الحوكمة والمؤسسية (١٥٪)", "٦٠.٠٠٪", "٩.٠٠ من ١٥ نقطة | ١٠٠٪ توطين وتطبيق قيود", 9.8)
    ]
    for title, pct, meta, left in bsc_cards:
        tb = s3.shapes.add_textbox(PInches(left), PInches(1.6), PInches(2.7), PInches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(12)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = pct
        p2.font.size = PPt(28)
        p2.font.bold = True
        p2.font.color.rgb = C_PRIMARY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = meta
        p3.font.size = PPt(10)
        p3.font.color.rgb = C_TEXT
        p3.alignment = PP_ALIGN.CENTER

    tb_score = s3.shapes.add_textbox(PInches(0.8), PInches(5.8), PInches(11.733), PInches(1.2))
    tf_score = tb_score.text_frame
    p_sc = tf_score.paragraphs[0]
    p_sc.text = "نسبة الإنجاز الاستراتيجي الإجمالية الموزونة: ٣٢.٧٢٪ | التقييم العام: يحتاج إلى تحسين جذري وإعادة ضبط مسار"
    p_sc.font.size = PPt(15)
    p_sc.font.bold = True
    p_sc.font.color.rgb = C_PRIMARY
    p_sc.alignment = PP_ALIGN.CENTER

    # Slide 4: Plan vs Actual Table
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "مصفوفة مطابقة الخطة الاستراتيجية بالمنجز الفعلي", "مقارنة أبرز المؤشرات المعتمدة لعام ٢٠٢٦م")
    table_shape = s4.shapes.add_table(7, 6, PInches(0.8), PInches(1.5), PInches(11.733), PInches(5.2))
    table = table_shape.table
    table_headers = ["المؤشر المعتمد", "المستهدف السنوي", "مستهدف H1", "المنجز الفعلي H1", "نسبة الإنجاز", "الحالة"]
    for j, h in enumerate(table_headers):
        table.cell(0, j).text = h
        table.cell(0, j).text_frame.paragraphs[0].font.bold = True
        table.cell(0, j).text_frame.paragraphs[0].font.size = PPt(11)
        table.cell(0, j).text_frame.paragraphs[0].font.color.rgb = C_WHITE
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = C_PRIMARY
    
    t_rows = [
        ("إجمالي الإيرادات الكلية", "٦,٨٤٦,٠٠٠ ر.س", "٣,٤٢٣,٠٠٠ ر.س", "٥٨٢,١٦٧ ر.س", "١٧.٠١٪", "متأخر حرِج"),
        ("المساعدات العلاجية المباشرة", "١,٥٠٠,٠٠٠ ر.س", "٧٥٠,٠٠٠ ر.س", "٢٠٨,٦٠٥ ر.س", "٢٧.٨١٪", "متأخر حرِج"),
        ("عدد المستفيدين المخدومين", "٣٦,٦٠٦ مستفيد", "١٨,٣٠٣ مستفيد", "٧ مستفيدين", "٠.٠٣٨٪", "متعثر تماماً"),
        ("الاستشارات الطبية", "١,٢٠٠ استشارة", "٦٠٠ استشارة", "٠ استشارة", "٠.٠٠٪", "لم يبدأ"),
        ("الشراكات الصحية الفاعلة", "٩ شراكات", "٩ شراكات", "٩ شراكات", "١٠٠.٠٠٪", "مكتمل"),
        ("توطين الكادر البشري", "١٠٠٪", "١٠٠٪", "١٠٠٪", "١٠٠.٠٠٪", "مكتمل")
    ]
    for i, r in enumerate(t_rows):
        for j, val in enumerate(r):
            table.cell(i+1, j).text = val
            table.cell(i+1, j).text_frame.paragraphs[0].font.size = PPt(10.5)

    # Slide 5: Critical Gap Analysis
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "تحليل الفجوات الاستراتيجية وإدارة المخاطر", "أبرز التحديات ونقاط الضعف المرصودة بالتدقيق الأدائي")
    gaps_ppt = [
        ("١. انحسار نطاق المستفيدين", "خدمة ٧ مرضى فقط من أصل ٣٦,٦٠٦ مستهدفين (-٩٩.٩٦٪ فجوة) بسبب حصر الصرف في الجراحات المعقدة وإيقاف العيادات الوقائية.", 0.8, 1.6),
        ("٢. تشدد لائحة المساعدات", "رفض ٦٦.٧٪ من الحالات المتقدمة (١٤ حالة) مما تسبب في عجز صرف موازنة العلاج (صرف ٢٠٨ ألف من ٧٥٠ ألف مخصصة).", 6.8, 1.6),
        ("٣. مخاطر تركز الإيرادات", "اعتماد بنسبة ٤٣٪ على متبرع فرد واحد (٢٥٠ ألف ريال) وتدني نسبة نجاح طلبات المنح إلى ٧.٤٪ فقط.", 0.8, 4.3),
        ("٤. تراجع النشاط التطوعي", "تنفيذ ٤ فرص تطوعية فقط دون توثيق للساعات مقابل مستهدف ٣,٠٠٠ ساعة بقيمة ٢٠٢.٥ ألف ريال.", 6.8, 4.3)
    ]
    for title, desc, left, top in gaps_ppt:
        tb = s5.shapes.add_textbox(PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_TEXT
        p2.alignment = PP_ALIGN.RIGHT

    # Slide 6: Roadmap
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "خارطة طريق التصحيح والتوصيات الاستراتيجية (H2 2026)", "حزمة الإجراءات التصحيحية لاستعادة الأثر وتحقيق الاستدامة")
    recs_ppt = [
        ("١. تعديل لائحة المساعدات فوراً", "وضع استثناءات إنسانية مرنة لرفع نسبة قبول الحالات وصرف المتبقي من الموازنة (٥٤١ ألف ريال).", 0.8, 1.6),
        ("٢. تشغيل برامج الاستشارات والفحوصات", "إطلاق الاستشارات الهاتفية والقوافل الميدانية لخدمة آلاف المستفيدين بتكلفة تشغيلية منخفضة.", 6.8, 1.6),
        ("٣. اعتماد موازنة الحوكمة ومنصة نوى", "التعاقد مع الفريق الاستشاري (١٥-٢١ ألف ريال) لإنهاء متطلبات الامتثال وفتح المنح الكبرى.", 0.8, 4.3),
        ("٤. إطلاق «بطاقة طبيبي» وتوحيد الخطة", "استثمار شبكة المستشفيات لإصدار بطاقة مزايا صحية ومواءمة مستهدفات الخطة الاستراتيجية.", 6.8, 4.3)
    ]
    for title, desc, left, top in recs_ppt:
        tb = s6.shapes.add_textbox(PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_TEXT
        p2.alignment = PP_ALIGN.RIGHT

    prs.save(pptx_path)
    print(f"Saved PPTX: {pptx_path}")

# Output paths
word_out = os.path.join(v2_dir, "تقرير_جمعية_طبيبي_النصف_سنوي_٢٠٢٦_النسخة_التنفيذية.docx")
pptx_out = os.path.join(v2_dir, "عرض_تقديمي_جمعية_طبيبي_٢٠٢٦_النسخة_التنفيذية.pptx")

build_word_document(word_out)
build_powerpoint(pptx_out)

shutil.copy2(word_out, os.path.join(v1_dir, "تقرير_جمعية_طبيبي_النصف_سنوي_٢٠٢٦_النسخة_التنفيذية.docx"))
shutil.copy2(pptx_out, os.path.join(v1_dir, "عرض_تقديمي_جمعية_طبيبي_٢٠٢٦_النسخة_التنفيذية.pptx"))

print("=== Word and PPTX successfully generated and synced ===")

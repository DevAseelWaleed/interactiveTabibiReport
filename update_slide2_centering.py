# -*- coding: utf-8 -*-
import os, sys, shutil

sys.stdout.reconfigure(encoding='utf-8')
base_dir = r"e:\Work\زبون تقرير نصف سنوي طبيبي"
v2_dir = os.path.join(base_dir, "التقرير_الاحترافي_المطور")
v1_dir = os.path.join(base_dir, "التقرير_الجديد")
images_dir = os.path.join(v2_dir, "assets", "images")

# 1. Update generate_web_slides.py
web_slides_path = os.path.join(base_dir, "generate_web_slides.py")
with open(web_slides_path, "r", encoding="utf-8") as f:
    web_code = f.read()

# Update Slide 2 styling in web slides to ensure center alignment of images and cards
slide2_old = """        <!-- SLIDE 2: Leadership -->
        <div class="slide" data-slide="2">
            <div class="slide-header">
                <div>
                    <div class="slide-eyebrow">الرؤية والتمكين الوطني</div>
                    <div class="slide-title">القيادة الرشيدة ومجلس الإدارة</div>
                </div>
                <div class="slide-logo"><i class="fas fa-heart-pulse"></i> طبيبي</div>
            </div>
            <div class="grid-3" style="margin-top:10px;">
                <!-- Crown Prince -->
                <div class="card-box" style="text-align:center;">
                    <img src="assets/images/crown_prince.jpg" alt="ولي العهد" class="royal-portrait">
                    <h4 style="color:var(--primary); font-size:1.1rem;">صاحب السمو الملكي</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:8px;">الأمير محمد بن سلمان بن عبدالعزيز</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main);">«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة.»</p>
                </div>
                <!-- King Salman -->
                <div class="card-box" style="text-align:center; border:2px solid var(--secondary); background:#FFFDF9;">
                    <img src="assets/images/king_salman.jpg" alt="خادم الحرمين الشريفين" class="royal-portrait" style="width:110px; height:110px; border-width:3px;">
                    <h4 style="color:var(--primary); font-size:1.15rem;">خادم الحرمين الشريفين</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:8px;">الملك سلمان بن عبدالعزيز آل سعود</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main);">«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، ومؤسساتها الخيرية.»</p>
                </div>
                <!-- Prince Salman bin Sultan -->
                <div class="card-box" style="text-align:center;">
                    <img src="assets/images/prince_salman.jpg" alt="أمير منطقة المدينة المنورة" class="royal-portrait">
                    <h4 style="color:var(--primary); font-size:1.1rem;">صاحب السمو الملكي</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:8px;">الأمير سلمان بن سلطان بن عبدالعزيز</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main);">«نسعد بالإنجازات التي حققتها الجمعيات الأهلية كشريك استراتيجي في جودة الحياة.»</p>
                </div>
            </div>
            <div style="font-size:0.85rem; color:var(--text-muted); text-align:center;">
                رئيس مجلس الإدارة: أ.د. منصور بن محمد النزهة | إشراف المركز الوطني لتنمية القطاع غير الربحي
            </div>
        </div>"""

slide2_new = """        <!-- SLIDE 2: Leadership -->
        <div class="slide" data-slide="2">
            <div class="slide-header">
                <div>
                    <div class="slide-eyebrow">الرؤية والتمكين الوطني</div>
                    <div class="slide-title">القيادة الرشيدة ومجلس الإدارة</div>
                </div>
                <div class="slide-logo"><i class="fas fa-heart-pulse"></i> طبيبي</div>
            </div>
            <div class="grid-3" style="margin-top:10px; align-items:stretch;">
                <!-- Crown Prince (Right) -->
                <div class="card-box" style="text-align:center; display:flex; flex-direction:column; align-items:center; justify-content:flex-start;">
                    <div style="display:flex; justify-content:center; align-items:center; width:100%; margin-bottom:12px;">
                        <img src="assets/images/crown_prince.jpg" alt="ولي العهد" class="royal-portrait">
                    </div>
                    <h4 style="color:var(--primary); font-size:1.15rem; margin-bottom:3px;">صاحب السمو الملكي</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:10px;">الأمير محمد بن سلمان بن عبدالعزيز</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main); text-align:center;">«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة.»</p>
                </div>

                <!-- King Salman (Center / Middle) -->
                <div class="card-box" style="text-align:center; border:2px solid var(--secondary); background:#FFFDF9; display:flex; flex-direction:column; align-items:center; justify-content:flex-start; box-shadow:0 8px 25px rgba(201, 169, 110, 0.18);">
                    <div style="display:flex; justify-content:center; align-items:center; width:100%; margin-bottom:12px;">
                        <img src="assets/images/king_salman.jpg" alt="خادم الحرمين الشريفين" class="royal-portrait" style="width:120px; height:120px; border-width:3.5px;">
                    </div>
                    <h4 style="color:var(--primary); font-size:1.2rem; margin-bottom:3px;">خادم الحرمين الشريفين</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:10px;">الملك سلمان بن عبدالعزيز آل سعود</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main); text-align:center;">«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، ومؤسساتها الخيرية.»</p>
                </div>

                <!-- Prince Salman bin Sultan (Left) -->
                <div class="card-box" style="text-align:center; display:flex; flex-direction:column; align-items:center; justify-content:flex-start;">
                    <div style="display:flex; justify-content:center; align-items:center; width:100%; margin-bottom:12px;">
                        <img src="assets/images/prince_salman.jpg" alt="أمير منطقة المدينة المنورة" class="royal-portrait">
                    </div>
                    <h4 style="color:var(--primary); font-size:1.15rem; margin-bottom:3px;">صاحب السمو الملكي</h4>
                    <p style="font-size:0.85rem; font-weight:700; color:var(--text-muted); margin-bottom:10px;">الأمير سلمان بن سلطان بن عبدالعزيز</p>
                    <p style="font-size:0.85rem; font-style:italic; line-height:1.6; color:var(--text-main); text-align:center;">«نسعد بالإنجازات التي حققتها الجمعيات الأهلية كشريك استراتيجي في جودة الحياة.»</p>
                </div>
            </div>
            <div style="font-size:0.85rem; color:var(--text-muted); text-align:center; margin-top:8px;">
                رئيس مجلس الإدارة: أ.د. منصور بن محمد النزهة | إشراف المركز الوطني لتنمية القطاع غير الربحي
            </div>
        </div>"""

if slide2_old in web_code:
    web_code = web_code.replace(slide2_old, slide2_new)
    with open(web_slides_path, "w", encoding="utf-8") as f:
        f.write(web_code)
    print("Updated generate_web_slides.py with centered Slide 2 images.")

# Rebuild Web Slides
os.system(f'py -3 "{web_slides_path}"')

# 2. Rebuild PPTX with Centered Images in Slide 2
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

def rebuild_pptx(pptx_path):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]

    C_PRIMARY = PRGBColor(107, 29, 58)
    C_SECONDARY = PRGBColor(201, 169, 110)
    C_WHITE = PRGBColor(255, 255, 255)
    C_TEXT = PRGBColor(45, 45, 45)

    def add_header(slide, title, subtitle):
        tb = slide.shapes.add_textbox(PInches(0.8), PInches(0.4), PInches(11.733), PInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(22)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = PPt(12)
        p2.font.color.rgb = C_SECONDARY
        p2.alignment = PP_ALIGN.RIGHT

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank_layout)
    tb = s1.shapes.add_textbox(PInches(1.0), PInches(1.5), PInches(11.333), PInches(4.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "جمعية طبيبي الأهلية بالمدينة المنورة"
    p.font.size = PPt(20)
    p.font.color.rgb = C_SECONDARY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "التقرير النصف سنوي الشامل لعام ٢٠٢٦م"
    p2.font.size = PPt(32)
    p2.font.bold = True
    p2.font.color.rgb = C_PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "«مدعم بالتدقيق الأدائي ومطابقة مستهدفات الخطة الاستراتيجية والتشغيلية»"
    p3.font.size = PPt(16)
    p3.font.bold = True
    p3.font.color.rgb = C_PRIMARY
    p3.alignment = PP_ALIGN.CENTER

    # Slide 2: Royal Leadership with Centered Images
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "القيادة الرشيدة ومجلس الإدارة", "الالتزام بتوجيهات القيادة في تمكين القطاع الصحي غير الربحي")
    leaders_info = [
        ("صاحب السمو الملكي الأمير محمد بن سلمان", "ولي العهد رئيس مجلس الوزراء", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة.»", "crown_prince.jpg", 0.8),
        ("خادم الحرمين الشريفين الملك سلمان بن عبدالعزيز", "ملك المملكة العربية السعودية", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه ومؤسساتها الخيرية.»", "king_salman.jpg", 4.8),
        ("صاحب السمو الملكي الأمير سلمان بن سلطان", "أمير منطقة المدينة المنورة", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية كشريك استراتيجي في جودة الحياة.»", "prince_salman.jpg", 8.8)
    ]
    for name, title, quote, img_file, left in leaders_info:
        img_path = os.path.join(images_dir, img_file)
        if os.path.exists(img_path):
            # Centered Image inside the 3.7 inch card width
            s2.shapes.add_picture(img_path, PInches(left + 1.15), PInches(1.5), PInches(1.4), PInches(1.4))
        
        tb = s2.shapes.add_textbox(PInches(left), PInches(3.05), PInches(3.7), PInches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = PPt(13.5)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = PPt(10.5)
        p2.font.color.rgb = C_SECONDARY
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = f"\n{quote}"
        p3.font.size = PPt(10.5)
        p3.font.italic = True
        p3.font.color.rgb = C_TEXT
        p3.alignment = PP_ALIGN.CENTER

    # Slide 3: BSC
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "التدقيق الأدائي والتقييم الاستراتيجي الشامل (BSC)", "نتائج تقييم النصف الأول ٢٠٢٦م وفق مناظير بطاقة الأداء المتوازن")
    bsc_cards = [
        ("محور الأثر والبرامج الطبية (٤٠٪)", "١٣.٩٥٪", "٥.٥٨ من ٤٠ نقطة | خدمة ٧ مرضى من أصل ٣٦ ألف", 0.8),
        ("المحور المالي والموازنة (٣٠٪)", "٣٢.٩٦٪", "٩.٨٩ من ٣٠ نقطة | تحصيل ٥٨٢ ألف ر.س (+١٩٢٪)", 3.8),
        ("محور الشراكات والعمليات (١٥٪)", "٥٥.٠٠٪", "٨.٢٥ من ١٥ نقطة | ٩ شراكات مستشفيات فاعلة", 6.8),
        ("محور الحوكمة والمؤسسية (١٥٪)", "٦٠.٠٠٪", "٩.٠٠ من ١٥ نقطة | ١٠٠٪ توطين وتطبيق قيود", 9.8)
    ]
    for title, pct, meta, left in bsc_cards:
        tb = s3.shapes.add_textbox(PInches(left), PInches(1.6), PInches(2.7), PInches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(12)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = pct
        p2.font.size = PPt(28)
        p2.font.bold = True
        p2.font.color.rgb = C_PRIMARY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = meta
        p3.font.size = PPt(10)
        p3.font.color.rgb = C_TEXT
        p3.alignment = PP_ALIGN.CENTER

    tb_score = s3.shapes.add_textbox(PInches(0.8), PInches(5.8), PInches(11.733), PInches(1.2))
    tf_score = tb_score.text_frame
    p_sc = tf_score.paragraphs[0]
    p_sc.text = "نسبة الإنجاز الاستراتيجي الإجمالية الموزونة: ٣٢.٧٢٪ | التقييم العام: يحتاج إلى تحسين جذري وإعادة ضبط مسار"
    p_sc.font.size = PPt(15)
    p_sc.font.bold = True
    p_sc.font.color.rgb = C_PRIMARY
    p_sc.alignment = PP_ALIGN.CENTER

    # Slide 4: Matrix
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "مصفوفة مطابقة الخطة الاستراتيجية بالمنجز الفعلي", "مقارنة أبرز المؤشرات المعتمدة لعام ٢٠٢٦م")
    table_shape = s4.shapes.add_table(7, 6, PInches(0.8), PInches(1.5), PInches(11.733), PInches(5.2))
    table = table_shape.table
    table_headers = ["المؤشر المعتمد", "المستهدف السنوي", "مستهدف H1", "المنجز الفعلي H1", "نسبة الإنجاز", "الحالة"]
    for j, h in enumerate(table_headers):
        table.cell(0, j).text = h
        table.cell(0, j).text_frame.paragraphs[0].font.bold = True
        table.cell(0, j).text_frame.paragraphs[0].font.size = PPt(11)
        table.cell(0, j).text_frame.paragraphs[0].font.color.rgb = C_WHITE
        table.cell(0, j).fill.solid()
        table.cell(0, j).fill.fore_color.rgb = C_PRIMARY
    
    t_rows = [
        ("إجمالي الإيرادات الكلية", "٦,٨٤٦,٠٠٠ ر.س", "٣,٤٢٣,٠٠٠ ر.س", "٥٨٢,١٦٧ ر.س", "١٧.٠١٪", "متأخر حرِج"),
        ("المساعدات العلاجية المباشرة", "١,٥٠٠,٠٠٠ ر.س", "٧٥٠,٠٠٠ ر.س", "٢٠٨,٦٠٥ ر.س", "٢٧.٨١٪", "متأخر حرِج"),
        ("عدد المستفيدين المخدومين", "٣٦,٦٠٦ مستفيد", "١٨,٣٠٣ مستفيد", "٧ مستفيدين", "٠.٠٣٨٪", "متعثر تماماً"),
        ("الاستشارات الطبية", "١,٢٠٠ استشارة", "٦٠٠ استشارة", "٠ استشارة", "٠.٠٠٪", "لم يبدأ"),
        ("الشراكات الصحية الفاعلة", "٩ شراكات", "٩ شراكات", "٩ شراكات", "١٠٠.٠٠٪", "مكتمل"),
        ("توطين الكادر البشري", "١٠٠٪", "١٠٠٪", "١٠٠٪", "١٠٠.٠٠٪", "مكتمل")
    ]
    for i, r in enumerate(t_rows):
        for j, val in enumerate(r):
            table.cell(i+1, j).text = val
            table.cell(i+1, j).text_frame.paragraphs[0].font.size = PPt(10.5)

    # Slide 5: Gaps
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "تحليل الفجوات الاستراتيجية وإدارة المخاطر", "أبرز التحديات ونقاط الضعف المرصودة بالتدقيق الأدائي")
    gaps_ppt = [
        ("١. انحسار نطاق المستفيدين", "خدمة ٧ مرضى فقط من أصل ٣٦,٦٠٦ مستهدفين (-٩٩.٩٦٪ فجوة) بسبب حصر الصرف في الجراحات المعقدة وإيقاف العيادات الوقائية.", 0.8, 1.6),
        ("٢. تشدد لائحة المساعدات", "رفض ٦٦.٧٪ من الحالات المتقدمة (١٤ حالة) مما تسبب في عجز صرف موازنة العلاج (صرف ٢٠٨ ألف من ٧٥٠ ألف مخصصة).", 6.8, 1.6),
        ("٣. مخاطر تركز الإيرادات", "اعتماد بنسبة ٤٣٪ على متبرع فرد واحد (٢٥٠ ألف ريال) وتدني نسبة نجاح طلبات المنح إلى ٧.٤٪ فقط.", 0.8, 4.3),
        ("٤. تراجع النشاط التطوعي", "تنفيذ ٤ فرص تطوعية فقط دون توثيق للساعات مقابل مستهدف ٣,٠٠٠ ساعة بقيمة ٢٠٢.٥ ألف ريال.", 6.8, 4.3)
    ]
    for title, desc, left, top in gaps_ppt:
        tb = s5.shapes.add_textbox(PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_TEXT
        p2.alignment = PP_ALIGN.RIGHT

    # Slide 6: Roadmap
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "خارطة طريق التصحيح والتوصيات الاستراتيجية (H2 2026)", "حزمة الإجراءات التصحيحية لاستعادة الأثر وتحقيق الاستدامة")
    recs_ppt = [
        ("١. تعديل لائحة المساعدات فوراً", "وضع استثناءات إنسانية مرنة لرفع نسبة قبول الحالات وصرف المتبقي من الموازنة (٥٤١ ألف ريال).", 0.8, 1.6),
        ("٢. تشغيل برامج الاستشارات والفحوصات", "إطلاق الاستشارات الهاتفية والقوافل الميدانية لخدمة آلاف المستفيدين بتكلفة تشغيلية منخفضة.", 6.8, 1.6),
        ("٣. اعتماد موازنة الحوكمة ومنصة نوى", "التعاقد مع الفريق الاستشاري (١٥-٢١ ألف ريال) لإنهاء متطلبات الامتثال وفتح المنح الكبرى.", 0.8, 4.3),
        ("٤. إطلاق «بطاقة طبيبي» وتوحيد الخطة", "استثمار شبكة المستشفيات لإصدار بطاقة مزايا صحية ومواءمة مستهدفات الخطة الاستراتيجية.", 6.8, 4.3)
    ]
    for title, desc, left, top in recs_ppt:
        tb = s6.shapes.add_textbox(PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_TEXT
        p2.alignment = PP_ALIGN.RIGHT

    prs.save(pptx_path)
    print(f"Saved PPTX with centered images: {pptx_path}")

pptx_out = os.path.join(v2_dir, "عرض_تقديمي_جمعية_طبيبي_٢٠٢٦_النسخة_التنفيذية.pptx")
rebuild_pptx(pptx_out)
shutil.copy2(pptx_out, os.path.join(v1_dir, "عرض_تقديمي_جمعية_طبيبي_٢٠٢٦_النسخة_التنفيذية.pptx"))

print("All slide images centered successfully across Web and PPTX.")

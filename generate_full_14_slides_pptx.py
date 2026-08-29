# -*- coding: utf-8 -*-
"""
Generate the complete 14-slide executive presentation for Tabibi Civil Association H1 2026:
File: عرض_تقديمي_جمعية_طبيبي_النصف_سنوي_٢٠٢٦.pptx
With centered royal portrait images, visual layout cards, and verified audited performance data.
"""
import os, sys, shutil
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

base_dir = r"e:\Work\زبون تقرير نصف سنوي طبيبي"
v2_dir = os.path.join(base_dir, "التقرير_الاحترافي_المطور")
v1_dir = os.path.join(base_dir, "التقرير_الجديد")
images_dir = os.path.join(v2_dir, "assets", "images")

# Primary Brand Colors
C_PRIMARY = PRGBColor(107, 29, 58)      # #6B1D3A Maroon/Burgundy
C_SECONDARY = PRGBColor(201, 169, 110)  # #C9A96E Gold
C_DARK_BG = PRGBColor(36, 7, 19)        # #240713 Dark Accent
C_CARD_BG = PRGBColor(250, 248, 245)    # Warm White Card
C_BORDER = PRGBColor(230, 224, 214)     # Border Gray
C_WHITE = PRGBColor(255, 255, 255)
C_TEXT_DARK = PRGBColor(45, 45, 45)
C_TEXT_MUTED = PRGBColor(100, 100, 100)
C_SUCCESS = PRGBColor(46, 125, 50)
C_WARNING = PRGBColor(217, 119, 6)
C_DANGER = PRGBColor(198, 40, 40)

def create_presentation():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_header(slide, title, subtitle, tag="طبيبي ٢٠٢٦م"):
        # Header banner text box
        tb = slide.shapes.add_textbox(PInches(0.8), PInches(0.35), PInches(11.733), PInches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = 'Arial'
        p_sub.font.size = PPt(11.5)
        p_sub.font.color.rgb = C_SECONDARY
        p_sub.font.bold = True
        p_sub.alignment = PP_ALIGN.RIGHT
        
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = PPt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = C_PRIMARY
        p_title.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: Cover Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_DARK_BG
    bg1.line.fill.background()

    tb1 = s1.shapes.add_textbox(PInches(1.0), PInches(1.2), PInches(11.333), PInches(5.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "جمعية طبيبي الأهلية بالمدينة المنورة"
    p.font.name = 'Arial'
    p.font.size = PPt(22)
    p.font.bold = True
    p.font.color.rgb = C_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "التقرير النصف سنوي الشامل لعام ٢٠٢٦م"
    p2.font.name = 'Arial'
    p2.font.size = PPt(34)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "«مدعم بالتدقيق الأدائي الشامل ومطابقة مستهدفات الخطة الاستراتيجية والتشغيلية»"
    p3.font.name = 'Arial'
    p3.font.size = PPt(16)
    p3.font.bold = True
    p3.font.color.rgb = C_SECONDARY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf1.add_paragraph()
    p4.text = "\n« ثـقـة  •  أثــر  •  اسـتـدامـة »"
    p4.font.name = 'Arial'
    p4.font.size = PPt(20)
    p4.font.bold = True
    p4.font.color.rgb = C_WHITE
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf1.add_paragraph()
    p5.text = "\nالفترة من ١ يناير إلى ٣٠ يونيو ٢٠٢٦م | ترخيص المركز الوطني لتنمية القطاع غير الربحي رقم: (١٠٠٠٧٣٠٧٠٠)"
    p5.font.name = 'Arial'
    p5.font.size = PPt(12.5)
    p5.font.color.rgb = C_SECONDARY
    p5.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: Royal Leadership (Images Centered)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "القيادة الرشيدة ومجلس الإدارة", "الرؤية والتمكين الوطني للقطاع الصحي غير الربحي")

    leaders_info = [
        ("صاحب السمو الملكي الأمير محمد بن سلمان", "ولي العهد رئيس مجلس الوزراء", "«نهدف للوصول إلى قطاع غير ربحي مهم، مبادر وداعم ومؤثر في التعليم والصحة والثقافة، وسنعتمد عليه بشكل رئيسي.»", "crown_prince.jpg", 0.8),
        ("خادم الحرمين الشريفين الملك سلمان بن عبدالعزيز", "ملك المملكة العربية السعودية", "«ما يميز هذه البلاد هو حرص قادتها على الخير والتشجيع عليه، ومؤسساتها الخيرية في مختلف المجالات.»", "king_salman.jpg", 4.8),
        ("صاحب السمو الملكي الأمير سلمان بن سلطان", "أمير منطقة المدينة المنورة", "«نسعد بالإنجازات التي حققتها الجمعيات الأهلية بالمنطقة كشريك استراتيجي في تحسين جودة الحياة.»", "prince_salman.jpg", 8.8)
    ]

    for name, title, quote, img_file, left in leaders_info:
        # Card background
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(1.5), PInches(3.7), PInches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_SECONDARY if "الملك سلمان" in name else C_BORDER
        card.line.width = PPt(2 if "الملك سلمان" in name else 1)

        # Centered Image inside the 3.7 inch card width
        img_path = os.path.join(images_dir, img_file)
        if os.path.exists(img_path):
            img_w = 1.45 if "الملك سلمان" in name else 1.35
            img_left = left + (3.7 - img_w) / 2.0
            s2.shapes.add_picture(img_path, PInches(img_left), PInches(1.65), PInches(img_w), PInches(img_w))

        # Text Box
        tb = s2.shapes.add_textbox(PInches(left + 0.15), PInches(3.15), PInches(3.4), PInches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = 'Arial'
        p.font.size = PPt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.name = 'Arial'
        p2.font.size = PPt(10.5)
        p2.font.color.rgb = C_SECONDARY
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = f"\n{quote}"
        p3.font.name = 'Arial'
        p3.font.size = PPt(10)
        p3.font.italic = True
        p3.font.color.rgb = C_TEXT_DARK
        p3.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: Executive Summary & BSC Strategic Audit Score
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "الملخص التنفيذي والتقييم الاستراتيجي الشامل (BSC)", "نتائج قياس الأداء للنصف الأول ٢٠٢٦م وفق بطاقة الأداء المتوازن")

    # 4 BSC Cards
    bsc_cards = [
        ("محور الأثر والبرامج الطبية (٤٠٪)", "١٣.٩٥٪", "٥.٥٨ من ٤٠ نقطة\nخدمة ٧ مرضى من أصل ٣٦ ألف مستهدف", C_DANGER, 0.8),
        ("المحور المالي والموازنة (٣٠٪)", "٦٣.٠١٪", "١٨.٩٠ من ٣٠ نقطة\nتحصيل ٥٨٢ ألف ر.س (٧٦.٢٥٪ من H1)", C_SUCCESS, 3.8),
        ("محور الشراكات والعمليات (١٥٪)", "٥٥.٠٠٪", "٨.٢٥ من ١٥ نقطة\nتفعيل ٩ شراكات صحية ومؤسسية فاعلة", PRGBColor(2, 136, 209), 6.8),
        ("محور الحوكمة والمؤسسية (١٥٪)", "٦٠.٠٠٪", "٩.٠٠ من ١٥ نقطة\n١٠٠٪ توطين وتطبيق نظام قيود السحابي", C_SUCCESS, 9.8)
    ]

    for title, pct, meta, color_accent, left in bsc_cards:
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(1.5), PInches(2.7), PInches(4.1))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = color_accent
        card.line.width = PPt(2)

        tb = s3.shapes.add_textbox(PInches(left + 0.15), PInches(1.65), PInches(2.4), PInches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = PPt(11)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = pct
        p2.font.name = 'Arial'
        p2.font.size = PPt(26)
        p2.font.bold = True
        p2.font.color.rgb = color_accent
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = meta
        p3.font.name = 'Arial'
        p3.font.size = PPt(9.5)
        p3.font.color.rgb = C_TEXT_DARK
        p3.alignment = PP_ALIGN.CENTER

    # Overall Score Banner
    score_banner = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(5.8), PInches(11.733), PInches(1.2))
    score_banner.fill.solid()
    score_banner.fill.fore_color.rgb = C_PRIMARY
    score_banner.line.color.rgb = C_SECONDARY
    score_banner.line.width = PPt(1.5)

    tb_s = s3.shapes.add_textbox(PInches(0.9), PInches(5.85), PInches(11.533), PInches(1.1))
    tf_s = tb_s.text_frame
    p_sc = tf_s.paragraphs[0]
    p_sc.text = "نسبة الإنجاز الاستراتيجي الإجمالية الموزونة: ٤١.٧٣٪  |  التقييم العام المعتمد: يحتاج إلى تحسين جذري وإعادة ضبط مسار"
    p_sc.font.name = 'Arial'
    p_sc.font.size = PPt(13.5)
    p_sc.font.bold = True
    p_sc.font.color.rgb = C_WHITE
    p_sc.alignment = PP_ALIGN.CENTER

    p_sub_sc = tf_s.add_paragraph()
    p_sub_sc.text = "حالة الأهداف: ٤ أهداف مكتملة  |  هدفان قيد التنفيذ  |  ٤ أهداف متأخرة  |  ٤ أهداف متعثرة تماماً أو لم تبدأ"
    p_sub_sc.font.name = 'Arial'
    p_sub_sc.font.size = PPt(11)
    p_sub_sc.font.color.rgb = C_SECONDARY
    p_sub_sc.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 4: Strategic Plan vs Actual Performance Matrix Table
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "مصفوفة مطابقة الخطة الاستراتيجية بالمنجز الفعلي", "مقارنة تفصيلية لمستهدفات عام ٢٠٢٦م المعتمدة وما تحقق على أرض الواقع")

    table_shape = s4.shapes.add_table(11, 6, PInches(0.8), PInches(1.5), PInches(11.733), PInches(5.4))
    table = table_shape.table
    table_headers = ["المؤشر الاستراتيجي المعتمد", "المستهدف السنوي", "مستهدف H1", "المنجز الفعلي H1", "نسبة الإنجاز H1", "الحالة المعتمدة"]
    
    for j, h in enumerate(table_headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.size = PPt(11)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        matrix_data = [
        ("الإيرادات المالية الكلية (المعدلة)", "١,٥٠٠,٠٠٠ ر.س", "٧٥٠,٠٠٠ ر.س", "٥٨٢,١٦٧.٥٢ ر.س", "٧٧.٦٢٪", "متقدم ومتميز"),
        ("الاستشارات الطبية والدوائية", "١,٢٠٠ استشارة", "٦٠٠ استشارة", "٠ استشارة", "٠.٠٠٪", "لم يبدأ"),
        ("الدراسات واستطلاعات الرأي", "٦ دراسات", "٣ دراسات", "٠ دراسة", "٠.٠٠٪", "لم يبدأ"),
        ("ساعات وقيمة العمل التطوعي", "٣,٠٠٠ س (٢٠٢ ألف)", "١,٥٠٠ س (١٠١ ألف)", "٤ فرص تطوعية", "غير مدققة", "متعثر"),
        ("عقد الشراكات الصحية الفاعلة", "٩ شراكات", "٩ شراكات", "٩ شراكات مفعلة", "١٠٠.٠٠٪", "مكتمل"),
        ("توطين الوظائف والكادر البشري", "١٠٠٪", "١٠٠٪", "١٠٠٪ (٣ موظفين)", "١٠٠.٠٠٪", "مكتمل"),
        ("تدريب وتأهيل الكادر الإداري", "٤ دورات", "٢ دورة", "٨ دورات تدريبية", "٤٠٠.٠٠٪", "متقدم"),
        ("التحول الرقمي والمحاسبي", "نظام سحابي", "نظام سحابي", "تم تشغيل قيود", "١٠٠.٠٠٪", "مكتمل"),
        ("معايير الحوكمة ومنصة نوى", "١٠٠٪", "بدء الملف", "قيد الاستعانة باستشاري", "جاري العمل", "قيد التنفيذ"),
        ("تنويع مصادر الدخل الذاتي", "٦ مصادر", "٦ مصادر", "٦ مصادر نشطة", "١٠٠.٠٠٪", "مكتمل")
    ]

    for i, row in enumerate(matrix_data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if i % 2 == 0 else C_WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = PPt(10)
            p.font.color.rgb = C_TEXT_DARK
            p.alignment = PP_ALIGN.CENTER
            if j == 4:
                p.font.bold = True
            elif j == 5:
                p.font.bold = True
                p.font.color.rgb = C_SUCCESS if "مكتمل" in val else (C_WARNING if "متأخر" in val else C_DANGER)

    # =========================================================================
    # SLIDE 5: Financial Performance & Revenue Comparison
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "الأداء المالي ونمو الإيرادات (H1 2026 vs H1 2025)", "تحليل مصادر الدخل ومقارنة الأداء المالي بالنصف المقابل من العام السابق")

    t5_shape = s5.shapes.add_table(8, 5, PInches(0.8), PInches(1.5), PInches(7.2), PInches(5.4))
    t5 = t5_shape.table
    t5_headers = ["بند الإيراد", "H1 2026 (ريال)", "H1 2025 (ريال)", "التغير (ريال)", "نسبة النمو"]
    for j, h in enumerate(t5_headers):
        cell = t5.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.size = PPt(10.5)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

    rev_rows = [
        ("الزكاة (أموال مقيدة)", "٧٠,٠٠٠", "٨٠,٠٠٠", "-١٠,٠٠٠", "-١٣٪"),
        ("علاج المرضى (مقيد)", "٧٥,٠٠٠", "٢٥,٠٠٠", "+٥٠,٠٠٠", "+٢٠٠٪"),
        ("المتجر الإلكتروني", "١٠,٤٦٩", "١٢٤", "+١٠,٣٤٥", "+٨,٣٤٣٪"),
        ("منصة تبرع الوطنية", "١,٢٠٣", "١٣,٧٨٦", "-١٢,٥٨٣", "-٩١٪"),
        ("تبرعات ودعم عام", "٤٠٧,٤٩٥", "٦٢,٥٦٤", "+٣٤٤,٩٣١", "+٥٥١٪"),
        ("رسوم العضوية", "١٨,٠٠٠", "١٨,٠٠٠", "٠", "٠٪"),
        ("الإجمالي العام للإيرادات", "٥٨٢,١٦٧.٥٢", "١٩٩,٤٧٤.٠٠", "+٣٨٢,٦٩٣.٥٢", "+١٩٢٪")
    ]
    for i, row in enumerate(rev_rows):
        for j, val in enumerate(row):
            cell = t5.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRGBColor(240, 235, 225) if i == 6 else (C_CARD_BG if i % 2 == 0 else C_WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = PPt(9.5)
            p.alignment = PP_ALIGN.CENTER
            if i == 6 or j == 4:
                p.font.bold = True
                p.font.color.rgb = C_SUCCESS if "+" in val else (C_DANGER if "-" in val else C_PRIMARY)

    # Right side Financial Insight Box
    c_fin = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(8.3), PInches(1.5), PInches(4.233), PInches(5.4))
    c_fin.fill.solid()
    c_fin.fill.fore_color.rgb = C_CARD_BG
    c_fin.line.color.rgb = C_SECONDARY
    c_fin.line.width = PPt(1.5)

    tb_f = s5.shapes.add_textbox(PInches(8.5), PInches(1.65), PInches(3.833), PInches(5.1))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "مؤشرات المركز المالي والسيولة:"
    p_f.font.name = 'Arial'
    p_f.font.size = PPt(13)
    p_f.font.bold = True
    p_f.font.color.rgb = C_PRIMARY

    p_f2 = tf_f.add_paragraph()
    p_f2.text = "\n• إجمالي النقدية بالبنوك: ١,٠٠١,٧٥٤ ريال\n  - البنك الأهلي: ٩٣٠,٧٠٢ ريال\n  - مصرف الراجحي: ٧١,٠٥٢ ريال\n\n• هيكل الأموال المتاحة:\n  - أموال مقيدة (زكاة وعلاج): ٣٦٧,٠٩٣ ريال\n  - أموال غير مقيدة: ٦٣٤,٦٦١ ريال\n\n• كفاية الاحتياطي النقدي: تغطية المصروفات التشغيلية لمدة ١٢ شهراً.\n\n• مخاطر التركز: تبرع واحد من فاعل خير بـ ٢٥٠ ألف ريال يمثل ٤٣٪ من إجمالي الدخل."
    p_f2.font.name = 'Arial'
    p_f2.font.size = PPt(10.5)
    p_f2.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 6: Budget Execution Table
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_header(s6, "تنفيذ الموازنة التقديرية لعام ٢٠٢٦م", "تحليل نسب الصرف الفعلي مقارنة بالمستهدف المعتمد بالموازنة التقديرية")

    t6_shape = s6.shapes.add_table(10, 5, PInches(0.8), PInches(1.5), PInches(11.733), PInches(5.4))
    t6 = t6_shape.table
    t6_headers = ["بند الموازنة التقديرية", "المستهدف التقديري (ريال)", "المحقق / المنصرف الفعلي", "نسبة الإنجاز", "التحليل المالي والدلالة"]
    for j, h in enumerate(t6_headers):
        cell = t6.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.size = PPt(10.5)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

    bgt_rows = [
        ("التبرعات والدعم العام (الإيرادات)", "١,٥٢٧,٠٠٠", "٥٨٢,١٦٧.٥٢", "٣٨.١٢٪", "تحقيق ٣٨٪ من الموازنة المستهدفة بنمو ١٩٢٪ عن ٢٠٢٥"),
        ("المساعدات العلاجية المباشرة للمرضى", "٧٥٠,٠٠٠", "٢٠٨,٦٠٥.٠٠", "٢٧.٨١٪", "صرف ٢٨٪ فقط لوجود تشدد في شروط قبول الحالات"),
        ("الرواتب والأجور المباشرة", "٤٧٢,٠٠٠", "١٤٤,٤٠٥.٠٠", "٣٠.٥٩٪", "انضباط عالي وصرف أجور ٣ موظفين ومحاسب متعاون"),
        ("المصروفات العمومية والتشغيلية", "١٤٢,٣٠٠", "١٠٩,٨٦٩.٠٠", "٧٧.٢١٪", "تشمل إيجار المقر والكهرباء والرسوم ونظام قيود"),
        ("التطوير المالي والإداري والحوكمة", "٤١,٠٠٠", "٠.٠٠", "٠.٠٠٪", "معلق لحين اعتماد الفريق الاستشاري للامتثال ونوى"),
        ("الحملة الإعلامية والتسويق", "٣٠,٠٠٠", "٠.٠٠", "٠.٠٠٪", "لم تُصرف مخصصات إعلامية مباشرة بالنصف الأول"),
        ("شراء الأصول والتجهيزات الثابتة", "١٩,٤٥٠", "١٥,٦٢٠.٨٠", "٨٠.٣١٪", "تأثيث المقر الجديد وشراء أجهزة حاسب وطابعات"),
        ("إجمالي موازنة الإيرادات التقديرية", "١,٥٢٧,٠٠٠", "٥٨٢,١٦٧.٥٢", "٣٨.١٢٪", "تحقيق ٧٦.٢٥٪ من مستهدف H1 (٥٨٢ ألف من ٧٦٣ ألف)"),
        ("إجمالي موازنة المصروفات التقديرية", "١,٤٥٤,٧٥٠", "٢٤٩,٢٧٤.٠٠", "١٧.١٣٪", "تنفيذ ٣٤.٢٧٪ من موازنة النصف التشغيلية بانضباط تام")
    ]
    for i, row in enumerate(bgt_rows):
        for j, val in enumerate(row):
            cell = t6.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRGBColor(240, 235, 225) if i == 7 else (C_CARD_BG if i % 2 == 0 else C_WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = PPt(9.5)
            p.alignment = PP_ALIGN.CENTER
            if i == 7 or j == 3:
                p.font.bold = True
                p.font.color.rgb = C_PRIMARY if i == 7 else (C_SUCCESS if "٨٠" in val or "٧٧" in val else C_TEXT_DARK)

    # =========================================================================
    # SLIDE 7: Operational & Administrative Expense Breakdown
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_header(s7, "المصروفات التشغيلية وترشيد النفقات", "تفصيل بنود المصاريف الإدارية والوفر السنوي المحقق بالنصف الأول")

    # 4 Cards of expenses highlights
    exp_highlights = [
        ("إيجار المقر الجديد", "٦٣,٣٣٣ ر.س", "سداد إيجار المقر الجديد (وفر سنوي ٢٥,٠٠٠ ريال مقارنة بالمقر السابق).", 0.8, 1.5),
        ("التأمينات الاجتماعية", "١٤,٧٦٨ ر.س", "سداد الاشتراكات الرسمية لكافة الموظفين السعوديين بانتظام تام.", 6.8, 1.5),
        ("أجور المتعاونين والتطوير", "١٦,٠٠٠ ر.س", "١٣ ألف أجور محاسب ومتعاونين + ٣ آلاف لتطوير الموقع الإلكتروني.", 0.8, 4.3),
        ("المحاسب القانوني والكهرباء", "٨,٤٦٧ ر.س", "٤,٦٠٠ أتعاب تدقيق محاسبي + ٣,٨٦٧ فواتير كهرباء للمقر الجديد.", 6.8, 4.3)
    ]
    for title, val, desc, left, top in exp_highlights:
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER

        tb = s7.shapes.add_textbox(PInches(left + 0.2), PInches(top + 0.2), PInches(5.3), PInches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{title} : {val}"
        p.font.name = 'Arial'
        p.font.size = PPt(14)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = 'Arial'
        p2.font.size = PPt(11)
        p2.font.color.rgb = C_TEXT_DARK
        p2.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 8: Medical Assistance & Supported Cases Table (7 Cases)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_header(s8, "البرامج الطبية والحالات المدعومة (برنامج جودة حياة)", "تفصيل الدعم العلاجي المباشر للحالات السبع المعتمدة خلال النصف الأول")

    t8_shape = s8.shapes.add_table(9, 6, PInches(0.8), PInches(1.5), PInches(11.733), PInches(5.4))
    t8 = t8_shape.table
    t8_headers = ["م", "اسم المستفيد", "التشخيص الطبي", "المستشفى المعالج", "مبلغ الدعم (ريال)", "نسبة التحسن والرضا"]
    for j, h in enumerate(t8_headers):
        cell = t8.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.size = PPt(10.5)
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

    cases_rows = [
        ("١", "فايز أحمد عبدالعزيز", "سرطان الدم (أورام دم معقدة)", "المستشفى السعودي الألماني", "١٥٠,٠٠٠.٠٠", "تحسن كامل (١٠٠٪)"),
        ("٢", "زينب عمر علي", "سرطان نخر العظم التخصصي", "المستشفى السعودي الألماني", "٣٠,٠٠٠.٠٠", "تحسن كامل (١٠٠٪)"),
        ("٣", "كندفة محمد عتبة", "تنويم وعلاج ورعاية ملاحظة", "مدينة الملك سلمان الطبية", "٧,٠٠٠.٠٠", "تحسن كامل (١٠٠٪)"),
        ("٤", "شوق حسن الأنور", "منظار جراحي استكشافي", "المستشفى السعودي الألماني", "٧,٠٠٠.٠٠", "تحسن كامل (١٠٠٪)"),
        ("٥", "سامية سليمان محمد", "استئصال كتلة من الصدر", "مستشفى المواساة بالمدينة", "٦,٣٥٠.٠٠", "نجاح العملية (١٠٠٪)"),
        ("٦", "زبيدة شمس الدين خاتم", "استئصال ورم بالقولون", "المستشفى السعودي الألماني", "٦,٣٣٠.٠٠", "تحسن كامل (١٠٠٪)"),
        ("٧", "محمد أحمد الشرفي", "أشعة رنين مغناطيسي دقيقة", "مستشفى المواساة بالمدينة", "١,٩٢٥.٠٠", "اكتمال الفحص (١٠٠٪)"),
        ("—", "الإجمالي العام للدعم الطبي", "٧ حالات علاجية تخصصية كبرى", "مستشفيات مرجعية بالمدينة", "٢٠٨,٦٠٥.٠٠", "متوسط الكلفة: ٢٩,٨٠١ ر.س")
    ]
    for i, row in enumerate(cases_rows):
        for j, val in enumerate(row):
            cell = t8.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRGBColor(240, 235, 225) if i == 7 else (C_CARD_BG if i % 2 == 0 else C_WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = PPt(9.5)
            p.alignment = PP_ALIGN.CENTER
            if i == 7 or j == 4:
                p.font.bold = True
                p.font.color.rgb = C_PRIMARY

    # =========================================================================
    # SLIDE 9: Rejected Cases Analysis & Governance
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_header(s9, "تحليل الحالات غير المقبولة (١٤ حالة) وحوكمة الصرف", "أسباب الرفض ومقترحات تعديل اللائحة لرفع نسبة قبول المرضى")

    rej_cards = [
        ("انتهاء صلاحية الإقامة (٧ حالات - ٥٠٪)", "الحالات: بسمة هارون، سيد الأمين، فريدة عظيم، عطور عباس، هاجر الصادق، عبدالله دياب، أحمد خير.\n• التوصية: وضع استثناءات للحالات الإنسانية الحرجة بموافقة المجلس.", 0.8, 1.5),
        ("تغطية كاملة من جمعية أخرى (حالتان - ١٤.٣٪)", "الحالات: هديباء الجهني (مياه بيضاء)، علي قايد (شرايين).\n• التوصية: تفعيل الربط الإلكتروني لتفادي الازدواجية وسرعة توجيه الدعم لمرضى آخرين.", 6.8, 1.5),
        ("أخطاء بالتقرير الطبي والتواريخ (حالتان - ١٤.٣٪)", "الحالات: ريم فواز (ورم ليفي)، جوهرة خان (أخطاء تواريخ).\n• التوصية: توجيه المستفيدين لتصحيح المستندات الطبية وإعادة الرفع.", 0.8, 4.3),
        ("تأمين طبي ساري / انتهاء تأشيرة (٣ حالات - ٢١.٤٪)", "الحالات: فؤاد لطف (تأمين طبي)، حمزة هندية (انتهاء زيارة)، مزاهر الهادي (قيد تسليم التعميد).\n• التوصية: تسريع مسارات دراسة الحالة واستكمال المتطلبات.", 6.8, 4.3)
    ]
    for title, desc, left, top in rej_cards:
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(top), PInches(5.7), PInches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_WARNING

        tb = s9.shapes.add_textbox(PInches(left + 0.2), PInches(top + 0.2), PInches(5.3), PInches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = PPt(12.5)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = 'Arial'
        p2.font.size = PPt(10)
        p2.font.color.rgb = C_TEXT_DARK
        p2.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 10: Health Partnerships Network (9 Partners)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_header(s10, "منظومة الشراكات الصحية والمؤسسية (٩ جهات)", "شبكة التحالفات الطبية لتغطية الفحوصات والعمليات الجراحية والتأهيل")

    partners_grid = [
        ("المستشفى السعودي الألماني", "علاج الأورام وسرطانات الدم والمناظير التخصصية", 0.8, 1.5),
        ("مستشفى المواساة بالمدينة", "العمليات الجراحية المتقدمة والأشعة والرنين", 4.8, 1.5),
        ("مدينة الملك سلمان الطبية", "الرعاية المرجعية التخصصية والتنويم الطبي", 8.8, 1.5),
        ("مستشفى د. حامد الأحمدي", "جراحات اليوم الواحد والعيادات الاستشارية", 0.8, 3.4),
        ("مستشفى المدينة الوطني", "خدمات الطوارئ والملاحظة والتحاليل الدقيقة", 4.8, 3.4),
        ("مستشفى المدينة الطبي العام", "الفحوصات العامة ورعاية الأمراض المزمنة", 8.8, 3.4),
        ("مستشفى واد الطبي", "علاج الإصابات الرياضية وجراحة العظام", 0.8, 5.3),
        ("شركة مداواة ورعاية الطبية", "توفير الأدوية والمستلزمات الطبية المنزلية", 4.8, 5.3),
        ("جمعية جَنَى لتأهيل المعاقات", "التأهيل الطبي والتكامل لرعاية ذوي الإعاقة", 8.8, 5.3)
    ]
    for name, scope, left, top in partners_grid:
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(top), PInches(3.7), PInches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER

        tb = s10.shapes.add_textbox(PInches(left + 0.15), PInches(top + 0.15), PInches(3.4), PInches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = 'Arial'
        p.font.size = PPt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = scope
        p2.font.name = 'Arial'
        p2.font.size = PPt(9.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 11: Human Resources & Capacity Building
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_header(s11, "رأس المال البشري والنمو المؤسسي", "الهيكل التنفيذي، التوطين بنسبة ١٠٠٪، وتأهيل وتدريب الكوادر")

    card_hr1 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(1.5), PInches(5.7), PInches(5.4))
    card_hr1.fill.solid()
    card_hr1.fill.fore_color.rgb = C_CARD_BG
    card_hr1.line.color.rgb = C_BORDER

    tb_hr1 = s11.shapes.add_textbox(PInches(1.0), PInches(1.7), PInches(5.3), PInches(5.0))
    tf_hr1 = tb_hr1.text_frame
    tf_hr1.word_wrap = True
    p = tf_hr1.paragraphs[0]
    p.text = "الهيكل الوظيفي والكادر التنفيذي:"
    p.font.name = 'Arial'
    p.font.size = PPt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p2 = tf_hr1.add_paragraph()
    p2.text = "\n• أ. بيان سعد المحمدي: المدير التنفيذي (موظف رسمي).\n• أ. غدير أحمد الحربي: المسؤول المالي والمشرفة على البرامج.\n• أ. طراد محمد سمان: سكرتير تنفيذي.\n• أ. محمد الحسن بشير: محاسب قانوني متعاون.\n• أ. فيصل الجهني: مسؤول علاقات عامة وإعلام (متعاون).\n\n• نسبة التوطين (السعودة): ١٠٠٪ في كافة المناصب الرسمية."
    p2.font.name = 'Arial'
    p2.font.size = PPt(11)
    p2.font.color.rgb = C_TEXT_DARK

    card_hr2 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(6.8), PInches(1.5), PInches(5.7), PInches(5.4))
    card_hr2.fill.solid()
    card_hr2.fill.fore_color.rgb = C_CARD_BG
    card_hr2.line.color.rgb = C_BORDER

    tb_hr2 = s11.shapes.add_textbox(PInches(7.0), PInches(1.7), PInches(5.3), PInches(5.0))
    tf_hr2 = tb_hr2.text_frame
    tf_hr2.word_wrap = True
    p = tf_hr2.paragraphs[0]
    p.text = "مؤشرات التأهيل والتدريب والتطوع:"
    p.font.name = 'Arial'
    p.font.size = PPt(14)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p2 = tf_hr2.add_paragraph()
    p2.text = "\n• ٨ دورات تدريبية متخصصة: استفاد منها موظفان بالكادر لرفع الكفاءة المالية والمؤسسية (مقابل ٠ دورات في ٢٠٢٥).\n\n• أنظمة الحضور والأرشفة: تطبيق نظام البصمة الإلكترونية وضبط الدوام والأرشفة الرقمية للوثائق.\n\n• مؤشر التطوع: تسجيل ٤ فرص تطوعية بالنصف الأول، مع التوصية بإعادة تفعيل التطوع الصحي التخصصي للوصول لمستهدف ٣,٠٠٠ ساعة تطوعية."
    p2.font.name = 'Arial'
    p2.font.size = PPt(11)
    p2.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 12: Resource Development & Grants Portfolio (27 Requests)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_header(s12, "تنمية الموارد ومسار المنح التمويلية (٢٧ جهة)", "حالة المنح المقبولة، الطلبات قيد المتابعة، والاعتذارات لربطها بالحوكمة")

    grants_cards = [
        ("منح مقبولة ومحققة (٤٠,٠٠٠ ريال)", "• مؤسسة إبراهيم العنقري: ٢٠,٠٠٠ ريال.\n• وقف الشيخ عبدالعزيز أبو زيد: ٢٠,٠٠٠ ريال.\n(تحققتا رغم عدم اكتمال الحوكمة والقوائم المالية حينها).", 0.8),
        ("طلبات قيد الدراسة والمتابعة (١١ جهة)", "• صندوق دعم الجمعيات (مشروع جودة حياة).\n• أوقاف الشيخ صالح الراجحي (دفء وغذاء).\n• بنك البلاد، بنك الرياض، و ٦ شركات لحجاج الداخل.", 4.8),
        ("اعتذارات بسبب الموازنة أو الشروط (٨ جهات)", "• انتهاء الموازنات: مؤسسة الماجد، مؤسسة الشاوي، شركة طيبة للاستثمار.\n• اشتراط الحوكمة: أوقاف الضحيان، مؤسسة طلال، مؤسسة المهيدب.", 8.8)
    ]
    for title, desc, left in grants_cards:
        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(left), PInches(1.5), PInches(3.7), PInches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_BORDER

        tb = s12.shapes.add_textbox(PInches(left + 0.15), PInches(1.7), PInches(3.4), PInches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = PPt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY
        p.alignment = PP_ALIGN.RIGHT

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = 'Arial'
        p2.font.size = PPt(10.5)
        p2.font.color.rgb = C_TEXT_DARK
        p2.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 13: Critical Gap Analysis & H2 2026 Roadmap
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_header(s13, "تحليل الفجوات الاستراتيجية وخارطة طريق التصحيح", "أبرز الفجوات المرصودة بالتدقيق وحزمة الإجراءات التنفيذية للنصف الثاني")

    # Left: 4 Gaps
    card_gaps = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.8), PInches(1.5), PInches(5.7), PInches(5.4))
    card_gaps.fill.solid()
    card_gaps.fill.fore_color.rgb = C_CARD_BG
    card_gaps.line.color.rgb = C_DANGER

    tb_g = s13.shapes.add_textbox(PInches(1.0), PInches(1.7), PInches(5.3), PInches(5.0))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True
    p = tf_g.paragraphs[0]
    p.text = "الفجوات الاستراتيجية المرصودة (Audit Gaps):"
    p.font.name = 'Arial'
    p.font.size = PPt(13.5)
    p.font.bold = True
    p.font.color.rgb = C_DANGER

    p2 = tf_g.add_paragraph()
    p2.text = "\n١. فجوة أعداد المستفيدين: خدمة ٧ مرضى فقط مقابل مستهدف ٣٦ ألف لحصر الصرف بالجراحات باهظة التكلفة.\n\n٢. تشدد لائحة المساعدات: رفض ٦٦.٧٪ من الحالات مما عطل صرف موازنة العلاج (صرف ٢٠٨ ألف من ٧٥٠ ألف).\n\n٣. تركز الإيرادات: ٤٣٪ من الدخل من متبرع فردي واحد مع تدني نسبة نجاح المنح (٧.٤٪).\n\n٤. تراجع التطوع: ٤ فرص تطوعية فقط دون توثيق للساعات."
    p2.font.name = 'Arial'
    p2.font.size = PPt(10.5)
    p2.font.color.rgb = C_TEXT_DARK

    # Right: Roadmap
    card_rec = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(6.8), PInches(1.5), PInches(5.7), PInches(5.4))
    card_rec.fill.solid()
    card_rec.fill.fore_color.rgb = C_CARD_BG
    card_rec.line.color.rgb = C_PRIMARY

    tb_r = s13.shapes.add_textbox(PInches(7.0), PInches(1.7), PInches(5.3), PInches(5.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "خارطة طريق التصحيح للنصف الثاني (H2 2026):"
    p.font.name = 'Arial'
    p.font.size = PPt(13.5)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p2 = tf_r.add_paragraph()
    p2.text = "\n١. تعديل لائحة المساعدات فوراً: إدخال استثناءات إنسانية مرنة لرفع القبول وصرف المتبقي من الموازنة (٥٤١ ألف ر.س).\n\n٢. تشغيل برامج الاستشارات والفحوصات: إطلاق «استشارات» و «أطمئن» لخدمة آلاف المستفيدين بتكلفة منخفضة.\n\n٣. اعتماد موازنة الحوكمة ونوى: التعاقد مع الفريق الاستشاري (١٥-٢١ ألف ريال) لإنهاء الامتثال وفتح المنح الكبرى.\n\n٤. إطلاق «بطاقة طبيبي» وتوحيد الخطة: إصدار بطاقة مزايا صحية مع المستشفيات الشريكة ومواءمة الخطة الاستراتيجية."
    p2.font.name = 'Arial'
    p2.font.size = PPt(10.5)
    p2.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 14: Closing & Official Contacts
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    bg14 = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    bg14.fill.solid()
    bg14.fill.fore_color.rgb = C_DARK_BG
    bg14.line.fill.background()

    tb14 = s14.shapes.add_textbox(PInches(1.0), PInches(1.5), PInches(11.333), PInches(4.8))
    tf14 = tb14.text_frame
    tf14.word_wrap = True

    p = tf14.paragraphs[0]
    p.text = "شكراً لثقتكم ودعمكم المستمر لرسالة جمعية طبيبي الأهلية"
    p.font.name = 'Arial'
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = C_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf14.add_paragraph()
    p2.text = "«نسعى لمواصلة المسيرة نحو قطاع صحي غير ربحي مستدام ومؤثر في طيبة الطيبة»"
    p2.font.name = 'Arial'
    p2.font.size = PPt(16)
    p2.font.italic = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf14.add_paragraph()
    p3.text = "\n\nهاتف: 00966555606347   |   البريد: tabibi2025med@gmail.com   |   الموقع: المدينة المنورة - حي الفتح"
    p3.font.name = 'Arial'
    p3.font.size = PPt(13.5)
    p3.font.color.rgb = C_SECONDARY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf14.add_paragraph()
    p4.text = "\nالمركز الوطني لتنمية القطاع غير الربحي - ترخيص رقم: (١٠٠٠٧٣٠٧٠٠)"
    p4.font.name = 'Arial'
    p4.font.size = PPt(11.5)
    p4.font.color.rgb = C_WHITE
    p4.alignment = PP_ALIGN.CENTER

    # Save to both target directories
    out_v2 = os.path.join(v2_dir, "عرض_تقديمي_جمعية_طبيبي_النصف_سنوي_٢٠٢٦.pptx")
    out_v1 = os.path.join(v1_dir, "عرض_تقديمي_جمعية_طبيبي_النصف_سنوي_٢٠٢٦.pptx")
    
    prs.save(out_v2)
    prs.save(out_v1)
    print(f"Generated complete 14-slide presentation successfully: {out_v2}")

create_presentation()
